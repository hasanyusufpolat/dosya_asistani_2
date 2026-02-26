"""
PROFESYONEL DOSYA ANALİZ MODÜLÜ
Yapay zeka tabanlı dosya analizi ve karar mekanizması
Her dosya türü için yapısal bütünlük, okunabilirlik ve profesyonellik değerlendirmesi
Gelişmiş metrikler, dil analizi ve kalite skorlaması
"""

import os
import re
import logging
import datetime
import hashlib
from typing import Dict, Tuple, Optional, List, Any
from dataclasses import dataclass, field
from enum import Enum
from collections import Counter

# Loglama ayarları
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('analyzer.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class AnalysisDecision(Enum):
    """Analiz karar tipleri"""
    DIRECT_CONVERT = "DOĞRUDAN_DÖNÜŞTÜR"
    SMART_EDIT = "AKILLI_DÜZENLE_GEREKLİ"
    QUALITY_NEEDED = "KALİTE_OPTİMİZASYONU_GEREKLİ"
    OCR_NEEDED = "OCR_GEREKLİ"

class DocumentComplexity(Enum):
    """Belge karmaşıklık seviyeleri"""
    VERY_LOW = "çok_düşük"
    LOW = "düşük"
    MEDIUM = "orta"
    HIGH = "yüksek"
    VERY_HIGH = "çok_yüksek"

class Language(Enum):
    """Dil seçenekleri"""
    TURKISH = "tr"
    ENGLISH = "en"
    GERMAN = "de"
    FRENCH = "fr"
    SPANISH = "es"
    ITALIAN = "it"
    RUSSIAN = "ru"
    ARABIC = "ar"
    MIXED = "karışık"
    UNKNOWN = "bilinmiyor"

@dataclass
class AnalysisMetrics:
    """Gelişmiş analiz metrikleri"""
    # Temel metrikler
    total_words: int = 0
    unique_words: int = 0
    total_sentences: int = 0
    total_paragraphs: int = 0
    total_pages: int = 0
    total_slides: int = 0
    
    # Ortalamalar
    avg_word_length: float = 0.0
    avg_sentence_length: float = 0.0
    avg_paragraph_length: float = 0.0
    
    # Okunabilirlik
    flesch_score: float = 0.0
    readability_level: str = ""
    
    # Yapısal metrikler
    heading_count: int = 0
    list_count: int = 0
    table_count: int = 0
    image_count: int = 0
    chart_count: int = 0
    
    # Kalite metrikleri
    consistency_score: int = 0
    formatting_score: int = 0
    language_score: int = 0
    overall_quality: int = 0
    
    # Dil metrikleri
    language: Language = Language.UNKNOWN
    has_turkish_chars: bool = False
    character_distribution: Dict[str, int] = field(default_factory=dict)
    
    # Karmaşıklık
    complexity: DocumentComplexity = DocumentComplexity.MEDIUM
    complexity_score: int = 0
    
    # Hatalar
    potential_errors: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)
    suggestions: List[str] = field(default_factory=list)

@dataclass
class AnalysisResult:
    """Analiz sonuç veri yapısı (gelişmiş)"""
    decision: AnalysisDecision
    confidence: float  # 0-100 arası güven skoru
    issues: List[str]
    suggestions: List[str]
    structure_score: int  # 0-100 arası yapısal puan
    readability_score: int  # 0-100 arası okunabilirlik puanı
    file_type: str
    details: Dict
    metrics: AnalysisMetrics = None
    processing_time: float = 0.0
    file_hash: str = ""
    file_size: int = 0

class FileAnalyzer:
    """Profesyonel Dosya Analizörü - Gelişmiş Versiyon"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': self.analyze_pdf,
            'docx': self.analyze_word,
            'doc': self.analyze_word,
            'xlsx': self.analyze_excel,
            'xls': self.analyze_excel,
            'pptx': self.analyze_powerpoint,
            'ppt': self.analyze_powerpoint,
            'txt': self.analyze_text,
            'md': self.analyze_markdown,
            'rtf': self.analyze_text,
            'png': self.analyze_image,
            'jpg': self.analyze_image,
            'jpeg': self.analyze_image,
            'gif': self.analyze_image,
            'bmp': self.analyze_image,
            'tiff': self.analyze_image
        }
        
        # Dil karakter setleri
        self.language_chars = {
            'tr': set('ğüşıöçĞÜŞİÖÇ'),
            'de': set('äöüßÄÖÜ'),
            'fr': set('éèêëàâçôùûÿœæ'),
            'es': set('ñáéíóúüÑÁÉÍÓÚÜ'),
            'it': set('àèéìíîòóùú'),
            'ru': set('абвгдеёжзийклмнопрстуфхцчшщъыьэюя'),
            'ar': set('ابتثجحخدذرزسشصضطظعغفقكلمنهوي')
        }
        
        # Okunabilirlik eşikleri
        self.readability_thresholds = {
            'çok_kolay': (90, 100),
            'kolay': (70, 89),
            'orta': (50, 69),
            'zor': (30, 49),
            'çok_zor': (0, 29)
        }
    
    def analyze(self, file_path: str) -> AnalysisResult:
        """
        Ana analiz fonksiyonu (gelişmiş)
        Dosya türünü otomatik algılar ve uygun analiz metodunu çağırır
        """
        import time
        start_time = time.time()
        
        try:
            # Dosya bilgileri
            file_ext = os.path.splitext(file_path)[1].lower().replace('.', '')
            file_size = os.path.getsize(file_path)
            file_hash = self._calculate_file_hash(file_path)
            
            if file_ext not in self.supported_formats:
                metrics = self._create_basic_metrics()
                return AnalysisResult(
                    decision=AnalysisDecision.DIRECT_CONVERT,
                    confidence=100,
                    issues=["Desteklenmeyen dosya formatı, varsayılan dönüşüm uygulanacak"],
                    suggestions=[],
                    structure_score=50,
                    readability_score=50,
                    file_type=file_ext,
                    details={"warning": "Desteklenmeyen format"},
                    metrics=metrics,
                    processing_time=time.time() - start_time,
                    file_hash=file_hash,
                    file_size=file_size
                )
            
            # Uygun analiz metodunu çağır
            analyzer = self.supported_formats[file_ext]
            result = analyzer(file_path)
            
            # Ortak alanları güncelle
            result.processing_time = time.time() - start_time
            result.file_hash = file_hash
            result.file_size = file_size
            
            # Kalite puanlarını hesapla
            if result.metrics:
                self._calculate_quality_scores(result)
            
            logger.info(f"✅ Analiz tamamlandı: {file_path} - Karar: {result.decision.value} - Güven: %{result.confidence}")
            return result
            
        except Exception as e:
            logger.error(f"❌ Analiz hatası: {e}")
            import traceback
            traceback.print_exc()
            
            metrics = self._create_basic_metrics()
            return AnalysisResult(
                decision=AnalysisDecision.DIRECT_CONVERT,
                confidence=60,
                issues=[f"Analiz sırasında hata: {str(e)}"],
                suggestions=[],
                structure_score=30,
                readability_score=30,
                file_type="unknown",
                details={"error": str(e)},
                metrics=metrics,
                processing_time=time.time() - start_time,
                file_hash="",
                file_size=0
            )
    
    def _calculate_file_hash(self, file_path: str) -> str:
        """Dosya hash'i hesapla"""
        try:
            with open(file_path, 'rb') as f:
                return hashlib.md5(f.read()).hexdigest()[:16]
        except:
            return ""
    
    def _create_basic_metrics(self) -> AnalysisMetrics:
        """Temel metrik nesnesi oluştur"""
        return AnalysisMetrics()
    
    def _detect_language(self, text: str) -> Language:
        """Metnin dilini tespit et"""
        if not text:
            return Language.UNKNOWN
        
        text_lower = text.lower()
        scores = {}
        
        for lang, chars in self.language_chars.items():
            score = sum(1 for c in text_lower if c in chars)
            if score > 0:
                scores[lang] = score
        
        if scores:
            # En yüksek skorlu dili bul
            main_lang = max(scores, key=scores.get)
            
            # Karışık dil kontrolü
            total = sum(scores.values())
            if len(scores) > 1 and max(scores.values()) / total < 0.7:
                return Language.MIXED
            
            lang_map = {
                'tr': Language.TURKISH,
                'en': Language.ENGLISH,
                'de': Language.GERMAN,
                'fr': Language.FRENCH,
                'es': Language.SPANISH,
                'it': Language.ITALIAN,
                'ru': Language.RUSSIAN,
                'ar': Language.ARABIC
            }
            return lang_map.get(main_lang, Language.UNKNOWN)
        
        # İngilizce varsayılan
        return Language.ENGLISH
    
    def _calculate_flesch_score(self, text: str) -> Tuple[float, str]:
        """
        Flesch okunabilirlik skorunu hesapla
        Türkçe ve İngilizce için uyarlanmış
        """
        if not text:
            return 0.0, "bilinmiyor"
        
        words = text.split()
        sentences = re.split(r'[.!?]+', text)
        sentences = [s for s in sentences if s.strip()]
        
        if not words or not sentences:
            return 0.0, "bilinmiyor"
        
        # Ortalama kelime sayısı
        avg_words_per_sentence = len(words) / len(sentences)
        
        # Ortalama hece sayısı (basitleştirilmiş)
        vowels = set('aeıioöuüAEIİOÖUÜ')
        total_syllables = sum(1 for word in words for char in word if char in vowels)
        avg_syllables_per_word = total_syllables / len(words) if words else 0
        
        # Flesch formülü
        flesch_score = 206.835 - 1.015 * avg_words_per_sentence - 84.6 * avg_syllables_per_word
        flesch_score = max(0, min(100, flesch_score))
        
        # Seviye belirle
        level = "bilinmiyor"
        for name, (low, high) in self.readability_thresholds.items():
            if low <= flesch_score <= high:
                level = name
                break
        
        return flesch_score, level
    
    def _calculate_complexity(self, metrics: AnalysisMetrics) -> DocumentComplexity:
        """Belge karmaşıklığını hesapla"""
        score = 0
        
        # Kelime çeşitliliği
        if metrics.total_words > 0:
            diversity = metrics.unique_words / metrics.total_words
            if diversity < 0.3:
                score += 30  # Çok tekrarlı
            elif diversity < 0.5:
                score += 20
            elif diversity < 0.7:
                score += 10
        
        # Cümle uzunluğu
        if metrics.avg_sentence_length > 30:
            score += 30
        elif metrics.avg_sentence_length > 20:
            score += 20
        elif metrics.avg_sentence_length > 15:
            score += 10
        
        # Yapısal elemanlar
        score += min(30, metrics.heading_count * 5)
        score += min(30, metrics.table_count * 10)
        score += min(20, metrics.chart_count * 10)
        
        metrics.complexity_score = min(100, score)
        
        # Karmaşıklık seviyesi
        if score < 20:
            return DocumentComplexity.VERY_LOW
        elif score < 40:
            return DocumentComplexity.LOW
        elif score < 60:
            return DocumentComplexity.MEDIUM
        elif score < 80:
            return DocumentComplexity.HIGH
        else:
            return DocumentComplexity.VERY_HIGH
    
    def _calculate_quality_scores(self, result: AnalysisResult):
        """Kalite puanlarını hesapla"""
        metrics = result.metrics
        
        # Tutarlılık puanı
        consistency = 100
        if len(metrics.warnings) > 0:
            consistency -= len(metrics.warnings) * 5
        if len(metrics.potential_errors) > 0:
            consistency -= len(metrics.potential_errors) * 10
        metrics.consistency_score = max(0, consistency)
        
        # Formatlama puanı
        formatting = 80
        if metrics.heading_count > 0:
            formatting += 10
        if metrics.list_count > 0:
            formatting += 5
        if metrics.table_count > 0:
            formatting += 5
        metrics.formatting_score = min(100, formatting)
        
        # Dil puanı
        if metrics.language != Language.UNKNOWN:
            metrics.language_score = 90
        else:
            metrics.language_score = 50
        
        # Genel kalite
        metrics.overall_quality = (
            metrics.consistency_score * 0.3 +
            metrics.formatting_score * 0.3 +
            metrics.language_score * 0.2 +
            result.structure_score * 0.1 +
            result.readability_score * 0.1
        )
    
    # ========== WORD ANALİZİ (GELİŞTİRİLMİŞ) ==========
    def analyze_word(self, file_path: str) -> AnalysisResult:
        """Word dosyası analizi (gelişmiş)"""
        try:
            from docx import Document
            from docx.oxml.ns import qn
            
            doc = Document(file_path)
            metrics = AnalysisMetrics()
            
            # Metin içeriğini topla
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            text = ' '.join(full_text)
            
            # Temel metrikler
            words = text.split()
            metrics.total_words = len(words)
            metrics.unique_words = len(set(words))
            
            sentences = re.split(r'[.!?]+', text)
            metrics.total_sentences = len([s for s in sentences if s.strip()])
            
            metrics.total_paragraphs = len([p for p in doc.paragraphs if p.text.strip()])
            
            # Ortalamalar
            if metrics.total_words > 0:
                metrics.avg_word_length = sum(len(w) for w in words) / metrics.total_words
            
            if metrics.total_sentences > 0:
                metrics.avg_sentence_length = metrics.total_words / metrics.total_sentences
            
            if metrics.total_paragraphs > 0:
                metrics.avg_paragraph_length = metrics.total_words / metrics.total_paragraphs
            
            # Dil tespiti
            metrics.language = self._detect_language(text)
            metrics.has_turkish_chars = any(c in self.language_chars['tr'] for c in text)
            
            # Okunabilirlik
            flesch, level = self._calculate_flesch_score(text)
            metrics.flesch_score = flesch
            metrics.readability_level = level
            
            # Yapısal analiz
            for para in doc.paragraphs:
                if para.style and 'heading' in para.style.name.lower():
                    metrics.heading_count += 1
            
            metrics.table_count = len(doc.tables)
            
            # Stil analizi
            styles = set()
            for para in doc.paragraphs:
                if para.style:
                    styles.add(para.style.name)
            
            # Sorun tespiti
            issues = []
            suggestions = []
            
            empty_paras = sum(1 for p in doc.paragraphs if not p.text.strip())
            if empty_paras > metrics.total_paragraphs * 0.2:
                issues.append(f"Çok fazla boş paragraf (%{(empty_paras/metrics.total_paragraphs*100):.1f})")
                suggestions.append("Boş paragrafları temizleyin")
                metrics.warnings.append("Boş paragraf fazlalığı")
            
            if metrics.heading_count == 0 and metrics.total_paragraphs > 10:
                issues.append("Dökümanda başlık yok")
                suggestions.append("Başlıklar ekleyin")
                metrics.warnings.append("Başlık yapısı zayıf")
            
            if len(styles) > 10:
                issues.append("Çok fazla stil kullanılmış")
                suggestions.append("Stilleri standardize edin")
            
            # Karar mekanizması
            structure_score = 80
            readability_score = 70
            
            if metrics.heading_count > 0:
                structure_score += 10
            if metrics.table_count > 0:
                structure_score += 5
            if metrics.avg_sentence_length > 25:
                readability_score -= 10
                issues.append("Cümleler ortalamadan uzun")
            
            structure_score = min(100, structure_score)
            readability_score = min(100, readability_score)
            
            # Karmaşıklık
            metrics.complexity = self._calculate_complexity(metrics)
            
            # Karar
            if structure_score >= 75 and readability_score >= 70:
                decision = AnalysisDecision.DIRECT_CONVERT
                confidence = (structure_score + readability_score) // 2
            else:
                decision = AnalysisDecision.SMART_EDIT
                confidence = 100 - ((100 - structure_score) + (100 - readability_score)) // 2
            
            return AnalysisResult(
                decision=decision,
                confidence=confidence,
                issues=issues,
                suggestions=suggestions,
                structure_score=structure_score,
                readability_score=readability_score,
                file_type="word",
                details={
                    "total_paragraphs": metrics.total_paragraphs,
                    "has_headings": metrics.heading_count > 0,
                    "has_tables": metrics.table_count > 0,
                    "language": metrics.language.value
                },
                metrics=metrics
            )
            
        except Exception as e:
            logger.error(f"❌ Word analiz hatası: {e}")
            return AnalysisResult(
                decision=AnalysisDecision.DIRECT_CONVERT,
                confidence=70,
                issues=[f"Word analiz hatası: {str(e)}"],
                suggestions=[],
                structure_score=50,
                readability_score=50,
                file_type="word",
                details={"error": str(e)},
                metrics=self._create_basic_metrics()
            )
    
    # ========== PDF ANALİZİ (GELİŞTİRİLMİŞ) ==========
    def analyze_pdf(self, file_path: str) -> AnalysisResult:
        """PDF dosyası analizi (gelişmiş)"""
        try:
            import PyPDF2
            
            metrics = AnalysisMetrics()
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                metrics.total_pages = len(pdf_reader.pages)
                full_text = []
                empty_pages = 0
                
                for page_num in range(metrics.total_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    
                    if text and text.strip():
                        full_text.append(text)
                    else:
                        empty_pages += 1
                        metrics.potential_errors.append(f"Sayfa {page_num + 1} metin içermiyor")
                
                text = ' '.join(full_text)
                
                # Temel metrikler
                words = text.split()
                metrics.total_words = len(words)
                metrics.unique_words = len(set(words))
                
                sentences = re.split(r'[.!?]+', text)
                metrics.total_sentences = len([s for s in sentences if s.strip()])
                
                # Ortalamalar
                if metrics.total_words > 0:
                    metrics.avg_word_length = sum(len(w) for w in words) / metrics.total_words
                
                if metrics.total_sentences > 0:
                    metrics.avg_sentence_length = metrics.total_words / metrics.total_sentences
                
                # OCR gereksinimi
                needs_ocr = empty_pages > metrics.total_pages * 0.3
                if needs_ocr:
                    metrics.warnings.append("PDF taranmış görüntü içerebilir")
                
                # Dil tespiti
                metrics.language = self._detect_language(text)
                
                # Okunabilirlik
                flesch, level = self._calculate_flesch_score(text)
                metrics.flesch_score = flesch
                metrics.readability_level = level
                
                # Sorun tespiti
                issues = []
                suggestions = []
                
                if empty_pages > metrics.total_pages * 0.3:
                    issues.append(f"Çok sayıda boş/az içerikli sayfa (%{(empty_pages/metrics.total_pages*100):.1f})")
                    suggestions.append("PDF'in taranmış olabileceğini kontrol edin")
                
                if needs_ocr:
                    issues.append("PDF taranmış görüntü içerebilir (OCR gerekli)")
                    suggestions.append("OCR ile metin çıkarımı yapılmalı")
                
                if metrics.avg_word_length < 3:
                    issues.append("Anormal kısa kelimeler var (OCR hatası olabilir)")
                    suggestions.append("Metin kalitesi düşük, düzenleme gerekli")
                
                # Karar mekanizması
                structure_score = 80
                readability_score = 70
                
                if empty_pages > 0:
                    structure_score -= empty_pages * 2
                if metrics.total_words > 1000:
                    readability_score += 10
                
                structure_score = max(30, min(100, structure_score))
                readability_score = max(30, min(100, readability_score))
                
                # Karmaşıklık
                metrics.complexity = self._calculate_complexity(metrics)
                
                if needs_ocr:
                    decision = AnalysisDecision.OCR_NEEDED
                    confidence = 60
                elif structure_score < 60:
                    decision = AnalysisDecision.SMART_EDIT
                    confidence = 100 - ((100 - structure_score) + (100 - readability_score)) // 2
                else:
                    decision = AnalysisDecision.DIRECT_CONVERT
                    confidence = (structure_score + readability_score) // 2
                
                return AnalysisResult(
                    decision=decision,
                    confidence=confidence,
                    issues=issues,
                    suggestions=suggestions,
                    structure_score=structure_score,
                    readability_score=readability_score,
                    file_type="pdf",
                    details={
                        "total_pages": metrics.total_pages,
                        "empty_pages": empty_pages,
                        "needs_ocr": needs_ocr,
                        "avg_word_length": metrics.avg_word_length
                    },
                    metrics=metrics
                )
                
        except Exception as e:
            logger.error(f"❌ PDF analiz hatası: {e}")
            return AnalysisResult(
                decision=AnalysisDecision.DIRECT_CONVERT,
                confidence=70,
                issues=[f"PDF analiz hatası: {str(e)}"],
                suggestions=[],
                structure_score=50,
                readability_score=50,
                file_type="pdf",
                details={"error": str(e)},
                metrics=self._create_basic_metrics()
            )
    
    # ========== EXCEL ANALİZİ (GELİŞTİRİLMİŞ) ==========
    def analyze_excel(self, file_path: str) -> AnalysisResult:
        """Excel dosyası analizi (gelişmiş)"""
        try:
            import pandas as pd
            
            metrics = AnalysisMetrics()
            excel_file = pd.ExcelFile(file_path)
            
            sheet_names = excel_file.sheet_names
            total_sheets = len(sheet_names)
            
            issues = []
            suggestions = []
            
            for sheet in sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet)
                
                if df.empty:
                    metrics.warnings.append(f"Sayfa '{sheet}' boş")
                
                metrics.table_count += 1
                metrics.total_words += df.size
            
            # Metrikler
            metrics.total_paragraphs = total_sheets
            
            # Sorun tespiti
            if total_sheets > 10:
                issues.append("Çok fazla sayfa var (10+)")
                suggestions.append("Sayfaları gruplandırın")
            
            if metrics.table_count > 20:
                issues.append("Çok sayıda tablo var")
                suggestions.append("Tabloları birleştirin veya sadeleştirin")
            
            # Karar mekanizması
            structure_score = 80
            readability_score = 75
            
            if total_sheets < 5:
                structure_score += 10
            if metrics.table_count > 50:
                structure_score -= 20
                readability_score -= 10
            
            structure_score = min(100, max(30, structure_score))
            readability_score = min(100, max(30, readability_score))
            
            # Karmaşıklık
            metrics.complexity = DocumentComplexity.HIGH if total_sheets > 10 else DocumentComplexity.MEDIUM
            
            # Karar
            if structure_score >= 70 and readability_score >= 70:
                decision = AnalysisDecision.DIRECT_CONVERT
                confidence = (structure_score + readability_score) // 2
            else:
                decision = AnalysisDecision.SMART_EDIT
                confidence = 100 - ((100 - structure_score) + (100 - readability_score)) // 2
            
            return AnalysisResult(
                decision=decision,
                confidence=confidence,
                issues=issues,
                suggestions=suggestions,
                structure_score=structure_score,
                readability_score=readability_score,
                file_type="excel",
                details={
                    "total_sheets": total_sheets,
                    "sheet_names": sheet_names[:5]
                },
                metrics=metrics
            )
            
        except Exception as e:
            logger.error(f"❌ Excel analiz hatası: {e}")
            return AnalysisResult(
                decision=AnalysisDecision.DIRECT_CONVERT,
                confidence=70,
                issues=[f"Excel analiz hatası: {str(e)}"],
                suggestions=[],
                structure_score=50,
                readability_score=50,
                file_type="excel",
                details={"error": str(e)},
                metrics=self._create_basic_metrics()
            )
    
    # ========== POWERPOINT ANALİZİ (GELİŞTİRİLMİŞ) ==========
    def analyze_powerpoint(self, file_path: str) -> AnalysisResult:
        """PowerPoint dosyası analizi (gelişmiş)"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            metrics = AnalysisMetrics()
            
            metrics.total_slides = len(prs.slides)
            empty_slides = 0
            text_heavy_slides = 0
            image_heavy_slides = 0
            
            for slide in prs.slides:
                has_text = False
                has_images = False
                text_count = 0
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        has_text = True
                        text_count += len(shape.text)
                        metrics.total_words += len(shape.text.split())
                    
                    if hasattr(shape, "image") or (hasattr(shape, "shape_type") and shape.shape_type == 13):
                        has_images = True
                        metrics.image_count += 1
                
                if not has_text and not has_images:
                    empty_slides += 1
                elif text_count > 500:
                    text_heavy_slides += 1
                elif has_images and not has_text:
                    image_heavy_slides += 1
            
            # Sorun tespiti
            issues = []
            suggestions = []
            
            if empty_slides > metrics.total_slides * 0.2:
                issues.append("Çok sayıda boş slayt var")
                suggestions.append("Boş slaytları temizleyin")
            
            if text_heavy_slides > metrics.total_slides * 0.5:
                issues.append("Slaytlar çok fazla metin içeriyor")
                suggestions.append("Metinleri birden çok slayta bölün")
            
            if image_heavy_slides > metrics.total_slides * 0.3:
                issues.append("Metinsiz görsel slaytlar var")
                suggestions.append("Görsellere açıklama ekleyin")
            
            # Karar mekanizması
            structure_score = 80
            readability_score = 75
            
            if empty_slides > 0:
                structure_score -= empty_slides * 2
            if metrics.image_count > metrics.total_slides:
                structure_score += 10
            
            structure_score = min(100, max(30, structure_score))
            readability_score = min(100, max(30, readability_score))
            
            # Karmaşıklık
            metrics.complexity = DocumentComplexity.MEDIUM
            
            # Karar
            if structure_score >= 70 and readability_score >= 70:
                decision = AnalysisDecision.DIRECT_CONVERT
                confidence = (structure_score + readability_score) // 2
            else:
                decision = AnalysisDecision.SMART_EDIT
                confidence = 100 - ((100 - structure_score) + (100 - readability_score)) // 2
            
            return AnalysisResult(
                decision=decision,
                confidence=confidence,
                issues=issues,
                suggestions=suggestions,
                structure_score=structure_score,
                readability_score=readability_score,
                file_type="powerpoint",
                details={
                    "total_slides": metrics.total_slides,
                    "empty_slides": empty_slides,
                    "image_count": metrics.image_count
                },
                metrics=metrics
            )
            
        except Exception as e:
            logger.error(f"❌ PowerPoint analiz hatası: {e}")
            return AnalysisResult(
                decision=AnalysisDecision.DIRECT_CONVERT,
                confidence=70,
                issues=[f"PowerPoint analiz hatası: {str(e)}"],
                suggestions=[],
                structure_score=50,
                readability_score=50,
                file_type="powerpoint",
                details={"error": str(e)},
                metrics=self._create_basic_metrics()
            )
    
    # ========== METİN ANALİZİ (GELİŞTİRİLMİŞ) ==========
    def analyze_text(self, file_path: str) -> AnalysisResult:
        """Düz metin dosyası analizi (gelişmiş)"""
        try:
            # Encoding tespiti
            encodings = ['utf-8', 'windows-1254', 'iso-8859-9', 'latin1']
            content = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                        content = f.read()
                    break
                except:
                    continue
            
            if content is None:
                raise Exception("Dosya okunamadı")
            
            metrics = AnalysisMetrics()
            
            # Temel metrikler
            lines = content.split('\n')
            metrics.total_paragraphs = len([l for l in lines if l.strip()])
            
            words = content.split()
            metrics.total_words = len(words)
            metrics.unique_words = len(set(words))
            
            sentences = re.split(r'[.!?]+', content)
            metrics.total_sentences = len([s for s in sentences if s.strip()])
            
            # Ortalamalar
            if metrics.total_words > 0:
                metrics.avg_word_length = sum(len(w) for w in words) / metrics.total_words
            
            if metrics.total_sentences > 0:
                metrics.avg_sentence_length = metrics.total_words / metrics.total_sentences
            
            if metrics.total_paragraphs > 0:
                metrics.avg_paragraph_length = metrics.total_words / metrics.total_paragraphs
            
            # Dil tespiti
            metrics.language = self._detect_language(content)
            metrics.has_turkish_chars = any(c in self.language_chars['tr'] for c in content)
            
            # Okunabilirlik
            flesch, level = self._calculate_flesch_score(content)
            metrics.flesch_score = flesch
            metrics.readability_level = level
            
            # Başlık tespiti
            for line in lines[:20]:
                if line.isupper() and len(line.split()) < 10:
                    metrics.heading_count += 1
            
            # Liste tespiti
            for line in lines:
                if line.strip().startswith(('•', '-', '*', '→', '✓')):
                    metrics.list_count += 1
            
            # Sorun tespiti
            issues = []
            suggestions = []
            
            empty_lines = sum(1 for l in lines if not l.strip())
            if empty_lines > len(lines) * 0.3:
                issues.append("Çok fazla boş satır var")
                suggestions.append("Boş satırları temizleyin")
            
            very_long_lines = sum(1 for l in lines if len(l) > 200)
            if very_long_lines > 0:
                issues.append("Çok uzun satırlar var")
                suggestions.append("Uzun satırları bölün")
            
            if metrics.avg_sentence_length > 30:
                issues.append("Cümleler çok uzun")
                suggestions.append("Cümleleri kısaltın")
            
            # Karar mekanizması
            structure_score = 70
            readability_score = 70
            
            if metrics.heading_count > 0:
                structure_score += 10
            if metrics.list_count > 0:
                structure_score += 5
            if metrics.avg_sentence_length < 20:
                readability_score += 10
            
            structure_score = min(100, structure_score)
            readability_score = min(100, readability_score)
            
            # Karmaşıklık
            metrics.complexity = self._calculate_complexity(metrics)
            
            # Karar
            if structure_score >= 70 and readability_score >= 70:
                decision = AnalysisDecision.DIRECT_CONVERT
                confidence = (structure_score + readability_score) // 2
            else:
                decision = AnalysisDecision.SMART_EDIT
                confidence = 100 - ((100 - structure_score) + (100 - readability_score)) // 2
            
            return AnalysisResult(
                decision=decision,
                confidence=confidence,
                issues=issues,
                suggestions=suggestions,
                structure_score=structure_score,
                readability_score=readability_score,
                file_type="text",
                details={
                    "total_lines": len(lines),
                    "empty_lines": empty_lines,
                    "has_headings": metrics.heading_count > 0,
                    "has_lists": metrics.list_count > 0
                },
                metrics=metrics
            )
            
        except Exception as e:
            logger.error(f"❌ Metin analiz hatası: {e}")
            return AnalysisResult(
                decision=AnalysisDecision.DIRECT_CONVERT,
                confidence=70,
                issues=[f"Metin analiz hatası: {str(e)}"],
                suggestions=[],
                structure_score=50,
                readability_score=50,
                file_type="text",
                details={"error": str(e)},
                metrics=self._create_basic_metrics()
            )
    
    # ========== MARKDOWN ANALİZİ ==========
    def analyze_markdown(self, file_path: str) -> AnalysisResult:
        """Markdown dosyası analizi"""
        # Markdown da metin gibi analiz edilebilir
        return self.analyze_text(file_path)
    
    # ========== GÖRSEL ANALİZİ (GELİŞTİRİLMİŞ) ==========
    def analyze_image(self, file_path: str) -> AnalysisResult:
        """Görsel dosyası analizi (gelişmiş)"""
        try:
            from PIL import Image, ImageStat
            
            image = Image.open(file_path)
            metrics = AnalysisMetrics()
            
            # Görsel metrikleri
            width, height = image.size
            resolution = width * height
            aspect_ratio = width / height if height > 0 else 1
            
            # Renk analizi
            is_grayscale = image.mode in ('L', 'LA')
            is_color = image.mode in ('RGB', 'RGBA', 'CMYK')
            
            # Parlaklık analizi
            if is_grayscale:
                stat = ImageStat.Stat(image)
                brightness = stat.mean[0] if stat.mean else 128
            else:
                image_gray = image.convert('L')
                stat = ImageStat.Stat(image_gray)
                brightness = stat.mean[0] if stat.mean else 128
            
            # Kalite metrikleri
            is_low_res = resolution < 100000  # 300x300 altı
            is_high_res = resolution > 2000000  # 2MP üstü
            is_dark = brightness < 50
            is_bright = brightness > 200
            
            # OCR gereksinimi
            needs_ocr = True  # Görseller her zaman OCR gerektirir
            
            # Sorun tespiti
            issues = []
            suggestions = []
            
            if is_low_res:
                issues.append("Düşük çözünürlüklü görsel")
                suggestions.append("Daha yüksek çözünürlük önerilir")
                metrics.warnings.append("Düşük çözünürlük")
            
            if is_dark:
                issues.append("Çok karanlık görsel")
                suggestions.append("Parlaklığı artırın")
            
            if is_bright:
                issues.append("Çok parlak görsel")
                suggestions.append("Kontrastı artırın")
            
            # Karar mekanizması
            structure_score = 50
            readability_score = 50
            
            if is_high_res:
                structure_score += 20
            if not is_low_res:
                readability_score += 20
            if 50 < brightness < 200:
                readability_score += 10
            
            structure_score = min(100, structure_score)
            readability_score = min(100, readability_score)
            
            # Karar
            decision = AnalysisDecision.OCR_NEEDED
            confidence = structure_score
            
            return AnalysisResult(
                decision=decision,
                confidence=confidence,
                issues=issues + ["Görselden metin çıkarımı (OCR) gerekiyor"],
                suggestions=suggestions + ["OCR işlemi uygulanacak", "Metin düzenlemesi yapılacak"],
                structure_score=structure_score,
                readability_score=readability_score,
                file_type="image",
                details={
                    "width": width,
                    "height": height,
                    "resolution": resolution,
                    "aspect_ratio": aspect_ratio,
                    "mode": image.mode,
                    "brightness": brightness,
                    "is_low_res": is_low_res,
                    "needs_ocr": needs_ocr
                },
                metrics=metrics
            )
            
        except Exception as e:
            logger.error(f"❌ Görsel analiz hatası: {e}")
            return AnalysisResult(
                decision=AnalysisDecision.OCR_NEEDED,
                confidence=50,
                issues=[f"Görsel analiz hatası: {str(e)}", "Varsayılan olarak OCR uygulanacak"],
                suggestions=[],
                structure_score=30,
                readability_score=30,
                file_type="image",
                details={"error": str(e)},
                metrics=self._create_basic_metrics()
            )


# ========== KULLANIM KOLAYLIĞI FONKSİYONLARI ==========

def analyze_file(file_path: str) -> Tuple[str, float, List[str], Dict]:
    """
    Dosya analizi yap ve karar döndür
    Returns: (karar, güven, sorunlar, detaylar)
    """
    analyzer = FileAnalyzer()
    result = analyzer.analyze(file_path)
    
    return (
        result.decision.value,
        result.confidence,
        result.issues,
        result.details
    )

def get_detailed_analysis(file_path: str) -> Dict[str, Any]:
    """
    Detaylı analiz raporu döndür
    """
    analyzer = FileAnalyzer()
    result = analyzer.analyze(file_path)
    
    return {
        'decision': result.decision.value,
        'confidence': result.confidence,
        'issues': result.issues,
        'suggestions': result.suggestions,
        'structure_score': result.structure_score,
        'readability_score': result.readability_score,
        'file_type': result.file_type,
        'metrics': {
            'total_words': result.metrics.total_words if result.metrics else 0,
            'unique_words': result.metrics.unique_words if result.metrics else 0,
            'total_sentences': result.metrics.total_sentences if result.metrics else 0,
            'avg_sentence_length': result.metrics.avg_sentence_length if result.metrics else 0,
            'language': result.metrics.language.value if result.metrics else 'unknown',
            'complexity': result.metrics.complexity.value if result.metrics else 'medium',
            'flesch_score': result.metrics.flesch_score if result.metrics else 0
        },
        'processing_time': result.processing_time,
        'file_hash': result.file_hash,
        'file_size': result.file_size
    }


# ========== TEST FONKSİYONU ==========
if __name__ == "__main__":
    print("🔧 Dosya Analiz Modülü Test Ediliyor...")
    print("=" * 60)
    
    # Test metni
    test_text = """
    BU BİR TEST BAŞLIĞIDIR
    
    Bu bir test paragrafıdır. İçinde birden fazla cümle bulunur. 
    Bu cümleler otomatik olarak analiz edilecek ve metrikler hesaplanacaktır.
    
    • Madde 1
    • Madde 2
    • Madde 3
    
    Bu da başka bir paragraf. İkinci paragraf olarak analiz edilecek.
    """
    
    # Geçici dosya oluştur
    test_file = "test_analysis.txt"
    with open(test_file, 'w', encoding='utf-8') as f:
        f.write(test_text)
    
    # Analiz yap
    analyzer = FileAnalyzer()
    result = analyzer.analyze(test_file)
    
    print(f"\n📊 Analiz Sonucu:")
    print(f"  • Karar: {result.decision.value}")
    print(f"  • Güven: %{result.confidence}")
    print(f"  • Yapısal Puan: {result.structure_score}")
    print(f"  • Okunabilirlik: {result.readability_score}")
    
    if result.metrics:
        print(f"\n📈 Metrikler:")
        print(f"  • Kelime Sayısı: {result.metrics.total_words}")
        print(f"  • Benzersiz Kelime: {result.metrics.unique_words}")
        print(f"  • Cümle Sayısı: {result.metrics.total_sentences}")
        print(f"  • Ort. Cümle Uzunluğu: {result.metrics.avg_sentence_length:.1f}")
        print(f"  • Dil: {result.metrics.language.value}")
        print(f"  • Karmaşıklık: {result.metrics.complexity.value}")
        print(f"  • Okunabilirlik: {result.metrics.readability_level}")
    
    if result.issues:
        print(f"\n⚠️ Sorunlar:")
        for issue in result.issues:
            print(f"  • {issue}")
    
    if result.suggestions:
        print(f"\n💡 Öneriler:")
        for suggestion in result.suggestions:
            print(f"  • {suggestion}")
    
    # Temizlik
    os.remove(test_file)
    
    print("=" * 60)
    print("✅ Dosya Analiz Modülü hazır!")