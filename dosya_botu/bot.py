"""
ANA BOT DOSYASI - PROFESYONEL VERSİYON (GÜNCELLENMİŞ)
Gelişmiş hata yönetimi, loglama ve optimizasyon
Tüm modüller entegre edilmiştir
Yapay zeka destekli analiz, isimlendirme, sınıflandırma, özetleme, doğrulama ve kalite optimizasyonu
"""

import sys
import os
import datetime
import sqlite3
import logging
import asyncio
import traceback
from typing import Optional, Dict, Any, List, Union
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import Application, CommandHandler, CallbackQueryHandler, MessageHandler, filters, ContextTypes

# Tek instance kontrolü için lock dosyası
LOCK_FILE = 'bot.lock'

def check_single_instance():
    """Bot'un tek bir instance çalıştığından emin ol (gelişmiş)"""
    try:
        if os.name == 'nt':  # Windows
            if os.path.exists(LOCK_FILE):
                try:
                    with open(LOCK_FILE, 'r') as f:
                        old_pid = f.read().strip()
                    # Eski işlem var mı kontrol et
                    import subprocess
                    result = subprocess.run(f'tasklist /FI "PID eq {old_pid}"', 
                                           shell=True, capture_output=True, text=True)
                    if str(old_pid) in result.stdout:
                        print(f"❌ Bot zaten çalışıyor (PID: {old_pid})")
                        print("💡 Eğer bot çalışmıyorsa, bot.lock dosyasını silin")
                        sys.exit(1)
                except:
                    pass
            
            with open(LOCK_FILE, 'w') as f:
                f.write(str(os.getpid()))
            print(f"✅ Bot başlatıldı (PID: {os.getpid()})")
            
        else:  # Linux/Mac
            try:
                import fcntl
                lock_file = open(LOCK_FILE, 'w')
                try:
                    fcntl.lockf(lock_file, fcntl.LOCK_EX | fcntl.LOCK_NB)
                    lock_file.write(str(os.getpid()))
                    lock_file.flush()
                    print(f"✅ Bot başlatıldı (PID: {os.getpid()})")
                except IOError:
                    print("❌ Bot zaten çalışıyor!")
                    print("💡 Eğer bot çalışmıyorsa, bot.lock dosyasını silin")
                    sys.exit(1)
            except ImportError:
                print("⚠️ fcntl modülü bulunamadı, instance kontrolü yapılamıyor")
                
    except Exception as e:
        print(f"⚠️ Instance kontrolü sırasında hata: {e}")
        print("💡 Bot çalıştırılıyor...")

# Kendi modüllerimiz
from config import *
import database as db
import converters
import utils
import analyzer
import naming
import classifier
import summarizer
import validator
import quality_optimizer
from payments import (
    show_packages, show_package_detail, start_payment,
    confirm_payment, approve_payment, reject_payment, 
    cancel_payment, back_to_main, init_payments_table
)

# Loglama ayarları
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('bot.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# ========== YARDIMCI FONKSİYONLAR ==========
def get_user_rights_direct(user_id: int) -> int:
    """Kullanıcının kalan hakkını doğrudan veritabanından al"""
    try:
        conn = sqlite3.connect('database/bot.db')
        cursor = conn.cursor()
        cursor.execute("SELECT remaining_rights FROM users WHERE user_id = ?", (user_id,))
        result = cursor.fetchone()
        conn.close()
        return result[0] if result else 0
    except Exception as e:
        logger.error(f"❌ Hak sorgulanırken hata: {e}")
        return 0

def extract_text_from_file(file_path: str, file_type: str) -> str:
    """Dosyadan metin çıkar (gelişmiş hata yönetimi)"""
    try:
        if file_type == 'WORD':
            try:
                from docx import Document
                doc = Document(file_path)
                return '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
            except Exception as e:
                logger.error(f"❌ Word okuma hatası: {e}")
                return f"[Word dosyası okunamadı: {e}]"
        
        elif file_type == 'PDF':
            try:
                import PyPDF2
                text = ""
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                return text
            except Exception as e:
                logger.error(f"❌ PDF okuma hatası: {e}")
                return f"[PDF dosyası okunamadı: {e}]"
        
        elif file_type == 'EXCEL':
            try:
                import pandas as pd
                df = pd.read_excel(file_path)
                return df.to_string()
            except Exception as e:
                logger.error(f"❌ Excel okuma hatası: {e}")
                return f"[Excel dosyası okunamadı: {e}]"
        
        elif file_type == 'POWERPOINT':
            try:
                from pptx import Presentation
                text = ""
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + "\n"
                return text
            except Exception as e:
                logger.error(f"❌ PowerPoint okuma hatası: {e}")
                return f"[PowerPoint dosyası okunamadı: {e}]"
        
        elif file_type == 'GORSEL':
            try:
                from PIL import Image
                import pytesseract
                
                # Tesseract yolunu kontrol et
                if os.name == 'nt':  # Windows
                    possible_paths = [
                        r'C:\Program Files\Tesseract-OCR\tesseract.exe',
                        r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
                    ]
                    for path in possible_paths:
                        if os.path.exists(path):
                            pytesseract.pytesseract.tesseract_cmd = path
                            break
                
                image = Image.open(file_path)
                return pytesseract.image_to_string(image, lang='tur+eng')
            except ImportError:
                logger.error("❌ pytesseract veya PIL yüklü değil")
                return "[OCR kütüphaneleri bulunamadı]"
            except Exception as e:
                logger.error(f"❌ Görsel OCR hatası: {e}")
                return f"[Görsel okunamadı: {e}]"
        
        else:
            # TXT veya diğer formatlar
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"❌ Metin dosyası okuma hatası: {e}")
                return f"[Dosya okunamadı: {e}]"
                
    except Exception as e:
        logger.error(f"❌ Metin çıkarma genel hatası: {e}")
        return f"[İşlem hatası: {e}]"

async def check_business_hours() -> bool:
    """
    Çalışma saatlerini kontrol et (7/24 aktif)
    """
    return True  # Her zaman True döndür, 7/24 çalışsın

async def get_business_hours_message() -> str:
    """
    Çalışma saatleri dışında gösterilecek mesaj (artık kullanılmıyor)
    """
    return ""

def get_quality_keyboard() -> InlineKeyboardMarkup:
    """Kalite seçim menüsü oluştur"""
    keyboard = [
        [
            InlineKeyboardButton("📄 Taslak", callback_data="quality_draft"),
            InlineKeyboardButton("📑 Standart", callback_data="quality_standard"),
        ],
        [
            InlineKeyboardButton("✨ Profesyonel", callback_data="quality_professional"),
            InlineKeyboardButton("💎 Premium", callback_data="quality_premium"),
        ],
        [InlineKeyboardButton("◀️ Geri", callback_data="back_to_convert")]
    ]
    return InlineKeyboardMarkup(keyboard)

# ========== BOT KOMUTLARI ==========
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Start komutu - Kullanıcıyı karşılar"""
    user = update.effective_user
    db.register_user(user)
    
    # Kalan hak kontrolü - doğrudan veritabanından
    remaining = get_user_rights_direct(user.id)
    
    keyboard = [[InlineKeyboardButton("👋 Merhaba", callback_data="merhaba")]]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    welcome_message = f"""🤖 **Dosya Asistanı'na hoş geldiniz!** 

━━━━━━━━━━━━━━━━━━━━━
👤 **Kullanıcı:** {user.first_name}
📦 **Kalan Hakkınız:** {remaining} Dosya

📁 **Desteklenen Formatlar:**
• PDF • Word • Excel • PowerPoint • Görsel

✨ **Akıllı Özellikler:**
• 📛 Akıllı İsimlendirme
• 📄 Belge Türü Tanıma
• 📋 Belge Özetleme
• ✅ Hata ve Eksik Kontrolü
• ⭐ Kalite Optimizasyonu (Taslak/Standart/Profesyonel/Premium)

⚡ **7/24 HİZMETİNİZDEYİZ!**

━━━━━━━━━━━━━━━━━━━━━
Başlamak için butona tıklayın."""
    
    await update.message.reply_text(
        welcome_message,
        reply_markup=reply_markup,
        parse_mode='Markdown'
    )
    
    logger.info(f"✅ Kullanıcı girişi: {user.id} - {user.first_name} - Kalan hak: {remaining}")

async def button_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Buton tıklamalarını yönet"""
    query = update.callback_query
    await query.answer()
    
    user_id = update.effective_user.id
    
    if query.data == "merhaba":
        keyboard = [
            [InlineKeyboardButton("📎 Dosya Yükle", callback_data="dosya_yukle")],
            [InlineKeyboardButton("💳 Paket Satın Al", callback_data="show_packages")],
            [InlineKeyboardButton("📊 Kalan Haklarım", callback_data="check_rights")],
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        await query.edit_message_text(
            text="📂 **Dosya Asistanı hazır**\n\nNe yapmak istersiniz?\n\n"
                 "📎 **Dosya Yükle** - Dönüşüm yapmak için\n"
                 "💳 **Paket Satın Al** - Yeni paket almak için\n"
                 "📊 **Kalan Haklarım** - Hak durumunuzu görmek için\n\n"
                 "✨ **Not:** Akıllı işlemler (isimlendirme, sınıflandırma, özetleme, vb.)\n"
                 "dosya yükledikten sonra otomatik olarak gösterilecektir.\n\n"
                 "⚡ **7/24 HİZMETİNİZDEYİZ**\n\n"
                 "Desteklenen dosya türleri:\n"
                 "• PDF (`.pdf`)\n"
                 "• Word (`.doc`, `.docx`)\n"
                 "• Excel (`.xls`, `.xlsx`)\n"
                 "• PowerPoint (`.ppt`, `.pptx`)\n"
                 "• Görsel (`.png`, `.jpg`, `.jpeg`)",
            reply_markup=reply_markup,
            parse_mode='Markdown'
        )
        logger.info(f"📋 Ana menü gösterildi: {user_id}")
    
    elif query.data == "dosya_yukle":
        await query.message.reply_text(
            "📎 **Dosya gönderme butonu**\n\n"
            "Lütfen aşağıdaki 📎 simgesine tıklayarak dosyanızı seçin ve gönderin."
        )
    
    elif query.data == "check_rights":
        remaining = get_user_rights_direct(user_id)
        stats = db.get_user_stats(user_id)
        
        if stats:
            message = f"""📊 **HAK DURUMUNUZ**

━━━━━━━━━━━━━━━━━━━━━
📦 **Kalan Hak:** `{remaining}` Dosya
✅ **Başarılı İşlem:** `{stats['success']}`
❌ **Başarısız İşlem:** `{stats['failed']}`
📈 **Toplam İşlem:** `{stats['total']}`
📅 **Bugünkü İşlem:** `{stats.get('today', 0)}`
📊 **Haftalık İşlem:** `{stats.get('weekly', 0)}`
📊 **Toplam Analiz:** `{stats.get('total_analysis', 0)}`
✨ **Akıllı Düzenleme:** `{stats.get('total_smart_edits', 0)}`
📛 **İsimlendirme:** `{stats.get('total_naming', 0)}`
📄 **Sınıflandırma:** `{stats.get('total_classification', 0)}`
📋 **Özetleme:** `{stats.get('total_summaries', 0)}`
✅ **Doğrulama:** `{stats.get('total_validations', 0)}`

━━━━━━━━━━━━━━━━━━━━━
💡 Yeni paket satın almak için butona tıklayın."""
            
            keyboard = [[InlineKeyboardButton("💳 Paket Satın Al", callback_data="show_packages")]]
            reply_markup = InlineKeyboardMarkup(keyboard)
            
            await query.message.reply_text(
                message,
                reply_markup=reply_markup,
                parse_mode='Markdown'
            )
        else:
            await query.message.reply_text("❌ Bilgilerinize ulaşılamadı.")

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Dosya gönderildiğinde çalışır"""
    user_id = update.effective_user.id
    username = update.effective_user.username or "isimsiz"
    
    # Hak kontrolü - doğrudan veritabanından
    remaining = get_user_rights_direct(user_id)
    logger.info(f"📥 Dosya yükleme: {username} - Kalan hak: {remaining}")
    
    if remaining <= 0:
        # Paket satın almak ister misiniz? sorusu
        keyboard = [
            [InlineKeyboardButton("✅ Evet, Paket Satın Al", callback_data="show_packages")],
            [InlineKeyboardButton("❌ Hayır, Teşekkürler", callback_data="merhaba")]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        await update.message.reply_text(
            "❌ **PAKET HAKKINIZ TÜKENMİŞTİR!**\n\n"
            "━━━━━━━━━━━━━━━━━━━━━\n"
            "📦 Dönüştürme işlemine devam etmek için yeni bir paket satın almanız gerekiyor.\n\n"
            "🎁 **SİZE ÖZEL İNDİRİMLİ PAKETLER:**\n"
            "• 🌟 Başlangıç Paketi: 5 Hak → 200 TL (300 TL)\n"
            "• 🚀 Gümüş Paket: 15 Hak → 500 TL (750 TL)\n"
            "• 💎 Elmas Paket: 30 Hak → 1000 TL (1400 TL) 🔥\n"
            "• 👑 Platin Paket: 50 Hak → 1500 TL (2000 TL)\n"
            "• 🏆 Elit Paket: 75 Hak → 2250 TL (3000 TL) 🔥\n\n"
            "━━━━━━━━━━━━━━━━━━━━━\n"
            "🔽 **Paket satın almak ister misiniz?**",
            reply_markup=reply_markup,
            parse_mode='Markdown'
        )
        return
    
    # Dosya bilgilerini al
    document = update.message.document
    file_name = document.file_name
    file_size = document.file_size
    
    # Dosya türünü belirle
    file_ext = os.path.splitext(file_name)[1].lower()
    
    if file_ext not in SUPPORTED_FORMATS:
        format_list = "\n".join([f"• {fmt}" for fmt in SUPPORTED_FORMATS.values()])
        await update.message.reply_text(
            f"❌ **Desteklenmeyen dosya türü!**\n\n"
            f"Lütfen şu formatlardan birini gönderin:\n{format_list}"
        )
        return
    
    file_type = SUPPORTED_FORMATS[file_ext]
    
    # temp klasörü kontrolü
    if not os.path.exists('temp'):
        os.makedirs('temp')
        logger.info("📁 temp klasörü oluşturuldu")
    
    # Dosyayı indir
    try:
        await update.message.reply_text(f"📥 **Dosya indiriliyor...**\nDosya: `{file_name}`")
        
        file = await context.bot.get_file(document.file_id)
        safe_name = utils.safe_filename(file_name)
        file_path = f"temp/{user_id}_{safe_name}"
        await file.download_to_drive(file_path)
        
        # Dosya boyutu kontrolü
        if file_size > MAX_FILE_SIZE:
            await update.message.reply_text(f"⚠️ **Dosya boyutu çok büyük!**\nMaksimum {MAX_FILE_SIZE/(1024*1024)} MB dosya gönderebilirsiniz.")
            os.remove(file_path)
            return
        
        await update.message.reply_text(f"✅ Dosya başarıyla indirildi.\nBoyut: `{utils.format_size(file_size)}`")
        logger.info(f"✅ Dosya indirildi: {file_name} - {utils.format_size(file_size)}")
        
    except Exception as e:
        logger.error(f"❌ Dosya indirilirken hata: {e}")
        await update.message.reply_text("❌ **Dosya indirilirken bir hata oluştu.**\nLütfen tekrar deneyin.")
        return
    
    # Dosya bilgisini kaydet
    context.user_data['current_file'] = file_path
    context.user_data['file_type'] = file_type
    context.user_data['file_name'] = file_name
    context.user_data['file_size'] = file_size
    
    # Dosya yüklendikten sonra işlem menüsünü göster
    await show_file_actions(update, context, file_path, file_type, file_name)

async def show_file_actions(update: Update, context: ContextTypes.DEFAULT_TYPE, file_path: str, file_type: str, file_name: str):
    """
    Dosya yüklendikten sonra yapılabilecek işlemleri göster
    - Önce dönüşüm seçenekleri
    - Sonra akıllı işlemler
    """
    
    # ÖNCE DÖNÜŞÜM SEÇENEKLERİ
    conversion_keyboard = []
    
    if file_type in CONVERSION_MAP:
        for target in CONVERSION_MAP[file_type]:
            display_name = DISPLAY_NAMES.get(target, target)
            callback_data = f"convert|{target}"
            conversion_keyboard.append([InlineKeyboardButton(f"🔄 {display_name}", callback_data=callback_data)])
    
    # SONRA AKILLI İŞLEMLER (dosya yüklendikten sonra)
    smart_keyboard = [
        [InlineKeyboardButton("━━━━━━━━━━━━━━━━━━━━━", callback_data="separator")],  # Ayraç
        [InlineKeyboardButton("✨ AKILLI İŞLEMLER", callback_data="smart_header")],  # Başlık
        [
            InlineKeyboardButton("📛 İsimlendir", callback_data="action_naming"),
            InlineKeyboardButton("📄 Sınıflandır", callback_data="action_classify"),
        ],
        [
            InlineKeyboardButton("✨ Akıllı Dönüşüm", callback_data="action_convert"),
            InlineKeyboardButton("📋 Özetle", callback_data="action_summarize"),
        ],
        [
            InlineKeyboardButton("✅ Hata Kontrolü", callback_data="action_validate"),
            InlineKeyboardButton("⭐ Kalite", callback_data="action_quality"),
        ],
        [InlineKeyboardButton("🔁 Tüm İşlemler (5 Hak)", callback_data="action_all")],
        [InlineKeyboardButton("◀️ Ana Menü", callback_data="merhaba")]
    ]
    
    # Klavyeleri birleştir
    keyboard = conversion_keyboard + smart_keyboard
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    # Dosya türüne göre açıklama metni
    conversion_options = ", ".join([DISPLAY_NAMES.get(t, t) for t in CONVERSION_MAP.get(file_type, [])])
    
    await update.message.reply_text(
        f"📂 **Dosya yüklendi:** `{file_name}`\n\n"
        f"📁 **Dosya türü:** {DISPLAY_NAMES.get(file_type, file_type)}\n"
        f"🔄 **Dönüştürülebilecek formatlar:** {conversion_options if conversion_options else 'Yok'}\n\n"
        f"**⚡ Ne yapmak istersiniz?**\n"
        f"• Yukarıdaki butonlardan birini seçin\n"
        f"• ✨ Akıllı işlemler için aşağı kaydırın",
        reply_markup=reply_markup,
        parse_mode='Markdown'
    )

# ========== KALİTE OPTİMİZASYONU HANDLER'LARI ==========
async def smart_quality_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Kalite optimizasyonu menüsü"""
    query = update.callback_query
    await query.answer()
    
    await query.edit_message_text(
        text="⭐ **KALİTE OPTİMİZASYONU**\n\n"
             "Lütfen istediğiniz kalite seviyesini seçin:\n\n"
             "📄 **Taslak** - Düşük kalite, hızlı işlem (1 hak)\n"
             "📑 **Standart** - Normal kalite (1 hak)\n"
             "✨ **Profesyonel** - Yüksek kalite, ofis standardı (1 hak)\n"
             "💎 **Premium** - Maksimum kalite, baskıya hazır (1 hak)",
        reply_markup=get_quality_keyboard(),
        parse_mode='Markdown'
    )

async def quality_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Kalite seçimi yapıldığında"""
    query = update.callback_query
    await query.answer()
    
    quality_map = {
        "quality_draft": "taslak",
        "quality_standard": "standart",
        "quality_professional": "profesyonel",
        "quality_premium": "premium"
    }
    
    quality = quality_map.get(query.data, "profesyonel")
    user_id = update.effective_user.id
    
    # Dosya bilgilerini al
    file_path = context.user_data.get('current_file')
    file_type = context.user_data.get('file_type')
    file_name = context.user_data.get('file_name')
    
    if not file_path or not os.path.exists(file_path):
        await query.edit_message_text("❌ **Dosya bulunamadı!**\nLütfen önce bir dosya yükleyin.")
        return
    
    await query.edit_message_text(f"⭐ **Kalite optimizasyonu yapılıyor...**\n\nSeçilen kalite: **{quality}**\n⏳ Lütfen bekleyin...")
    
    try:
        # Metin çıkar
        text_content = extract_text_from_file(file_path, file_type)
        
        # Kalite optimizasyonu yap
        optimizer = quality_optimizer.QualityOptimizer()
        quality_level = getattr(quality_optimizer.QualityLevel, quality.upper(), 
                               quality_optimizer.QualityLevel.PROFESSIONAL)
        optimizer.quality_level = quality_level
        
        success, new_path, results = optimizer.optimize_document(file_path, file_type, text_content)
        
        if success:
            # Hak tüket
            rights_success = db.decrease_rights(user_id, 'quality')
            
            if not rights_success:
                await query.edit_message_text("❌ Kalite optimizasyonu için yeterli hakkınız bulunmamaktadır.")
                return
            
            # Dosya yolunu güncelle
            context.user_data['current_file'] = new_path
            
            # Rapor göster
            report = optimizer.get_quality_report(results)
            
            message = f"✅ **Kalite Optimizasyonu Tamamlandı!**\n\n{report}"
            
            # Kalan hakkı göster
            remaining = get_user_rights_direct(user_id)
            message += f"\n🔁 **Kalan hak:** {remaining}"
            
            keyboard = [[InlineKeyboardButton("◀️ İşlemlere Dön", callback_data="merhaba")]]
            reply_markup = InlineKeyboardMarkup(keyboard)
            
            await query.edit_message_text(message, reply_markup=reply_markup, parse_mode='Markdown')
        else:
            await query.edit_message_text("❌ Kalite optimizasyonu başarısız oldu.")
            
    except Exception as e:
        logger.error(f"❌ Kalite optimizasyonu hatası: {e}")
        await query.edit_message_text(f"❌ Hata oluştu: {str(e)[:100]}")

async def back_to_convert_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Dönüşüm menüsüne geri dön"""
    query = update.callback_query
    await query.answer()
    
    file_type = context.user_data.get('file_type')
    
    if file_type:
        # Burada uygun işlem yapılabilir
        await query.edit_message_text("İşlem iptal edildi. Yeni dosya yükleyebilirsiniz.")
    else:
        # Ana menüye dön
        await query.edit_message_text("❌ Dosya bulunamadı. Ana menüye dönülüyor...")

# ========== AKILLI İŞLEM HANDLER'LARI ==========
async def smart_naming_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Akıllı isimlendirme butonu"""
    query = update.callback_query
    await query.answer()
    
    user_id = update.effective_user.id
    
    # Dosya bilgilerini al
    file_path = context.user_data.get('current_file')
    file_type = context.user_data.get('file_type')
    file_name = context.user_data.get('file_name')
    file_size = context.user_data.get('file_size', 0)
    
    if not file_path or not os.path.exists(file_path):
        await query.edit_message_text("❌ **Dosya bulunamadı!**\nLütfen önce bir dosya yükleyin.")
        return
    
    await handle_smart_naming(query, context, user_id, file_path, file_type, file_name, file_size)

async def smart_classify_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Belge türü tanıma butonu"""
    query = update.callback_query
    await query.answer()
    
    user_id = update.effective_user.id
    
    # Dosya bilgilerini al
    file_path = context.user_data.get('current_file')
    file_type = context.user_data.get('file_type')
    file_name = context.user_data.get('file_name')
    file_size = context.user_data.get('file_size', 0)
    
    if not file_path or not os.path.exists(file_path):
        await query.edit_message_text("❌ **Dosya bulunamadı!**\nLütfen önce bir dosya yükleyin.")
        return
    
    await handle_smart_classify(query, context, user_id, file_path, file_type, file_name, file_size)

async def smart_convert_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Akıllı dönüşüm butonu"""
    query = update.callback_query
    await query.answer()
    
    user_id = update.effective_user.id
    
    # Dosya bilgilerini al
    file_path = context.user_data.get('current_file')
    file_type = context.user_data.get('file_type')
    file_name = context.user_data.get('file_name')
    file_size = context.user_data.get('file_size', 0)
    
    if not file_path or not os.path.exists(file_path):
        await query.edit_message_text("❌ **Dosya bulunamadı!**\nLütfen önce bir dosya yükleyin.")
        return
    
    await handle_smart_convert(query, context, user_id, file_path, file_type, file_name, file_size)

async def smart_summarize_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Belge özetleme butonu"""
    query = update.callback_query
    await query.answer()
    
    user_id = update.effective_user.id
    
    # Dosya bilgilerini al
    file_path = context.user_data.get('current_file')
    file_type = context.user_data.get('file_type')
    file_name = context.user_data.get('file_name')
    file_size = context.user_data.get('file_size', 0)
    
    if not file_path or not os.path.exists(file_path):
        await query.edit_message_text("❌ **Dosya bulunamadı!**\nLütfen önce bir dosya yükleyin.")
        return
    
    await handle_smart_summarize(query, context, user_id, file_path, file_type, file_name, file_size)

async def smart_validate_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Hata ve eksik kontrolü butonu"""
    query = update.callback_query
    await query.answer()
    
    user_id = update.effective_user.id
    
    # Dosya bilgilerini al
    file_path = context.user_data.get('current_file')
    file_type = context.user_data.get('file_type')
    file_name = context.user_data.get('file_name')
    file_size = context.user_data.get('file_size', 0)
    
    if not file_path or not os.path.exists(file_path):
        await query.edit_message_text("❌ **Dosya bulunamadı!**\nLütfen önce bir dosya yükleyin.")
        return
    
    await handle_smart_validate(query, context, user_id, file_path, file_type, file_name, file_size)

async def smart_all_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Tüm akıllı işlemler butonu"""
    query = update.callback_query
    await query.answer()
    
    user_id = update.effective_user.id
    
    # Dosya bilgilerini al
    file_path = context.user_data.get('current_file')
    file_type = context.user_data.get('file_type')
    file_name = context.user_data.get('file_name')
    file_size = context.user_data.get('file_size', 0)
    
    if not file_path or not os.path.exists(file_path):
        await query.edit_message_text("❌ **Dosya bulunamadı!**\nLütfen önce bir dosya yükleyin.")
        return
    
    await handle_smart_all(query, context, user_id, file_path, file_type, file_name, file_size)

async def smart_action_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Akıllı işlem butonlarını yönet (dosya yüklendikten sonra)"""
    query = update.callback_query
    await query.answer()
    
    # Ayraç ve başlık butonlarını yoksay
    if query.data in ["separator", "smart_header"]:
        return
    
    action = query.data
    user_id = update.effective_user.id
    
    # Dosya bilgilerini al
    file_path = context.user_data.get('current_file')
    file_type = context.user_data.get('file_type')
    file_name = context.user_data.get('file_name')
    file_size = context.user_data.get('file_size', 0)
    
    if not file_path or not os.path.exists(file_path):
        await query.edit_message_text("❌ **Dosya bulunamadı!**\nLütfen dosyayı tekrar yükleyin.")
        return
    
    if action == "action_naming":
        await handle_smart_naming(query, context, user_id, file_path, file_type, file_name, file_size)
    
    elif action == "action_classify":
        await handle_smart_classify(query, context, user_id, file_path, file_type, file_name, file_size)
    
    elif action == "action_convert":
        await handle_smart_convert(query, context, user_id, file_path, file_type, file_name, file_size)
    
    elif action == "action_summarize":
        await handle_smart_summarize(query, context, user_id, file_path, file_type, file_name, file_size)
    
    elif action == "action_validate":
        await handle_smart_validate(query, context, user_id, file_path, file_type, file_name, file_size)
    
    elif action == "action_quality":
        await smart_quality_handler(update, context)
    
    elif action == "action_all":
        await handle_smart_all(query, context, user_id, file_path, file_type, file_name, file_size)

async def handle_smart_naming(query, context, user_id, file_path, file_type, file_name, file_size):
    """Akıllı isimlendirme işlemi"""
    await query.edit_message_text("📛 **Akıllı isimlendirme yapılıyor...**\n\n⏳ Lütfen bekleyin...")
    
    try:
        # Metin çıkar
        text_content = extract_text_from_file(file_path, file_type)
        
        # İsimlendirme yap
        success, new_path, naming_info = naming.smart_rename(file_path, text_content)
        
        if success:
            # Hak tüket
            rights_success = db.increase_naming_count(user_id)
            
            if not rights_success:
                await query.edit_message_text("❌ İsimlendirme için yeterli hakkınız bulunmamaktadır.")
                return
            
            # Kaydet
            db.save_naming_record(
                user_id=user_id,
                original_name=file_name,
                new_name=naming_info['new_name'],
                extracted_info=naming_info['extracted_info'],
                confidence=naming_info['confidence']
            )
            
            # Dosya yolunu güncelle
            context.user_data['current_file'] = new_path
            
            message = f"""✅ **Akıllı İsimlendirme Tamamlandı!**

━━━━━━━━━━━━━━━━━━━━━
📛 **Yeni isim:** `{naming_info['new_name']}`
📊 **Güven skoru:** %{naming_info['confidence']}

📋 **Çıkarılan bilgiler:**
"""
            for key, value in naming_info['extracted_info'].items():
                message += f"• {key}: {value}\n"
            
            message += f"\n📌 **1 hak tüketildi.**"
            
            # Kalan hakkı göster
            remaining = get_user_rights_direct(user_id)
            message += f"\n🔁 **Kalan hak:** {remaining}"
            
            keyboard = [[InlineKeyboardButton("◀️ İşlemlere Dön", callback_data="merhaba")]]
            reply_markup = InlineKeyboardMarkup(keyboard)
            
            await query.edit_message_text(message, reply_markup=reply_markup, parse_mode='Markdown')
            
        else:
            await query.edit_message_text("❌ İsimlendirme başarısız oldu.")
            
    except Exception as e:
        logger.error(f"❌ İsimlendirme hatası: {e}")
        await query.edit_message_text(f"❌ Hata oluştu: {str(e)[:100]}")

async def handle_smart_classify(query, context, user_id, file_path, file_type, file_name, file_size):
    """Belge türü tanıma işlemi"""
    await query.edit_message_text("📄 **Belge türü analiz ediliyor...**\n\n⏳ Lütfen bekleyin...")
    
    try:
        # Metin çıkar
        text_content = extract_text_from_file(file_path, file_type)
        
        # Sınıflandır
        class_result = classifier.classify_document(text_content, file_type)
        
        # Hak tüket
        rights_success = db.increase_classification_count(user_id)
        
        if not rights_success:
            await query.edit_message_text("❌ Sınıflandırma için yeterli hakkınız bulunmamaktadır.")
            return
        
        # Kaydet
        db.save_classification_record(
            user_id=user_id,
            file_name=file_name,
            document_type=class_result['document_type'],
            category=class_result['category'],
            confidence=class_result['confidence'],
            allowed_formats=class_result['allowed_formats'],
            extracted_fields=class_result['extracted_fields']
        )
        
        message = f"""✅ **Belge Türü Analizi Tamamlandı!**

━━━━━━━━━━━━━━━━━━━━━
📄 **Belge türü:** `{class_result['document_type']}`
📂 **Kategori:** `{class_result['category']}`
📊 **Güven skoru:** %{class_result['confidence']}

✅ **İzin verilen formatlar:** {', '.join(class_result['allowed_formats'])}

📋 **Çıkarılan alanlar:**
"""
        for key, value in class_result['extracted_fields'].items():
            message += f"• {key}: {value}\n"
        
        message += f"\n📌 **1 hak tüketildi.**"
        
        # Kalan hakkı göster
        remaining = get_user_rights_direct(user_id)
        message += f"\n🔁 **Kalan hak:** {remaining}"
        
        keyboard = [[InlineKeyboardButton("◀️ İşlemlere Dön", callback_data="merhaba")]]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        await query.edit_message_text(message, reply_markup=reply_markup, parse_mode='Markdown')
        
    except Exception as e:
        logger.error(f"❌ Sınıflandırma hatası: {e}")
        await query.edit_message_text(f"❌ Hata oluştu: {str(e)[:100]}")

async def handle_smart_convert(query, context, user_id, file_path, file_type, file_name, file_size):
    """Akıllı dönüşüm işlemi"""
    # Önce dönüşüm formatını sor
    keyboard = []
    options = CONVERSION_MAP.get(file_type, ['PDF'])
    
    for opt in options:
        display_name = DISPLAY_NAMES.get(opt, opt)
        callback_data = f"smart_convert_to|{opt}"
        keyboard.append([InlineKeyboardButton(display_name, callback_data=callback_data)])
    
    keyboard.append([InlineKeyboardButton("❌ İptal", callback_data="action_cancel")])
    
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await query.edit_message_text(
        f"📄 **Dosya türü:** {DISPLAY_NAMES.get(file_type, file_type)}\n\n"
        f"Lütfen dönüştürmek istediğiniz formatı seçin:",
        reply_markup=reply_markup,
        parse_mode='Markdown'
    )

# ========== YENİ EKLENEN - GELİŞMİŞ DÖNÜŞÜM HANDLER'ı ==========
async def smart_convert_to_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """
    Dönüşüm formatı seçildikten sonra akıllı dönüşüm yap - GELİŞMİŞ VERSİYON
    """
    query = update.callback_query
    await query.answer()
    
    target = query.data.replace('smart_convert_to|', '')
    user_id = update.effective_user.id
    
    # Dosya bilgilerini al
    file_path = context.user_data.get('current_file')
    file_type = context.user_data.get('file_type')
    file_name = context.user_data.get('file_name')
    file_size = context.user_data.get('file_size', 0)
    
    # DEBUG: Gelen değerleri logla
    logger.info(f"🔍 smart_convert_to_handler çağrıldı - Hedef: {target}, Kaynak tip: {file_type}")
    
    # Dosya kontrolü
    if not file_path or not os.path.exists(file_path):
        await query.edit_message_text("❌ **Dosya bulunamadı!**\nLütfen dosyayı tekrar yükleyin.")
        return
    
    # Hedef tip kontrolü
    if not target or target == "UNKNOWN":
        await query.edit_message_text("❌ **Geçersiz dönüşüm hedefi!**\nLütfen tekrar deneyin.")
        return
    
    # Çıktı dosyası adını oluştur
    base_name = os.path.splitext(file_name)[0]
    output_ext = EXTENSION_MAP.get(target, '.pdf')
    safe_name = utils.safe_filename(f"{user_id}_{base_name}_converted{output_ext}")
    output_path = f"temp/{safe_name}"
    
    # Kullanıcıya bilgi ver
    await query.edit_message_text(
        f"🔄 **Dönüşüm başlatılıyor...**\n\n"
        f"📁 Kaynak: `{file_name}`\n"
        f"🎯 Hedef: **{DISPLAY_NAMES.get(target, target)}**\n"
        f"⚙️ Kalite: Profesyonel\n\n"
        f"⏳ Bu işlem birkaç saniye sürebilir..."
    )
    
    try:
        # Dönüşümü yap
        start_time = datetime.datetime.now()
        
        success, out_path, conv_type, edit_summary, metrics = await converters.smart_convert_file(
            input_path=file_path,
            output_path=output_path,
            source_type=file_type,
            target_type=target,
            user_id=user_id,
            db_instance=db,
            quality="profesyonel"
        )
        
        processing_time = (datetime.datetime.now() - start_time).total_seconds()
        
        if success and os.path.exists(out_path):
            # Hak tüket
            rights_consumed = 1
            db.decrease_rights(user_id, conv_type)
            
            # Dönüşüm kaydını ekle
            db.save_conversion_record(
                user_id=user_id,
                file_name=file_name,
                file_size=file_size,
                source_format=file_type,
                target_format=target,
                conversion_type=conv_type,
                status='success',
                processing_time=processing_time,
                quality_score=metrics.quality_score if metrics else 95
            )
            
            # Kullanıcıya dönüştürülmüş dosyayı gönder
            with open(out_path, 'rb') as f:
                await query.message.reply_document(
                    document=f,
                    filename=f"{base_name}_converted{output_ext}",
                    caption=f"✅ **Dönüştürme tamamlandı!**"
                )
            
            # Başarı mesajı
            quality_text = f"%{metrics.quality_score}" if metrics and metrics.quality_score else "Yüksek"
            await query.message.reply_text(
                f"✅ **Dönüşüm başarılı!**\n"
                f"📁 {DISPLAY_NAMES.get(file_type, file_type)} → {DISPLAY_NAMES.get(target, target)}\n"
                f"⏱️ Süre: {processing_time:.1f}s\n"
                f"📊 Kalite: {quality_text}\n"
                f"📌 {rights_consumed} hak tüketildi."
            )
            
            # Rapor göster (varsa)
            if metrics and hasattr(metrics, 'warnings') and metrics.warnings:
                warning_text = "\n".join([f"⚠️ {w}" for w in metrics.warnings[:3]])
                await query.message.reply_text(f"📋 **Dönüşüm notları:**\n{warning_text}")
            
            # Kalan hakkı göster
            new_remaining = get_user_rights_direct(user_id)
            await query.message.reply_text(f"🔁 **Kalan hak:** {new_remaining}")
            
            # Geçici dosyaları temizle
            utils.clean_temp_files(user_id, file_path, out_path)
            
        else:
            # Başarısız dönüşüm
            error_msg = metrics.warnings[0] if metrics and metrics.warnings else "Bilinmeyen hata"
            
            db.increase_failed_count(user_id)
            db.save_conversion_record(
                user_id=user_id,
                file_name=file_name,
                file_size=file_size,
                source_format=file_type,
                target_format=target,
                conversion_type='direct',
                status='failed',
                processing_time=processing_time,
                error_message=error_msg
            )
            
            # Hata mesajı
            await query.message.reply_text(
                f"❌ **Dönüştürme başarısız!**\n\n"
                f"📁 Dosya: `{file_name}`\n"
                f"⚠️ Hata: {error_msg}\n\n"
                f"📞 Destek: @yusozone"
            )
            
            # Geçici dosyayı temizle (kaynak dosya silinmez)
            utils.clean_temp_files(user_id, file_path)
            
    except Exception as e:
        logger.error(f"❌ Akıllı dönüşüm hatası: {e}")
        traceback.print_exc()
        await query.message.reply_text(
            f"❌ **Dönüştürme sırasında hata oluştu!**\n\n"
            f"⚠️ Hata: {str(e)[:200]}\n\n"
            f"📞 Destek: @yusozone"
        )
        utils.clean_temp_files(user_id, file_path)

async def handle_smart_summarize(query, context, user_id, file_path, file_type, file_name, file_size):
    """Belge özetleme işlemi"""
    await query.edit_message_text("📋 **Belge özetleniyor...**\n\n⏳ Lütfen bekleyin...")
    
    try:
        # Metin çıkar
        text_content = extract_text_from_file(file_path, file_type)
        
        # Belge türünü belirle
        doc_type = file_type.lower()
        
        # Özetle
        summary_result = summarizer.summarize_document(text_content, doc_type)
        
        # Hak tüket
        rights_success = db.increase_summary_count(user_id)
        
        if not rights_success:
            await query.edit_message_text("❌ Özetleme için yeterli hakkınız bulunmamaktadır.")
            return
        
        # Kaydet
        db.save_summary_record(
            user_id=user_id,
            file_name=file_name,
            summary=summary_result['summary'],
            key_points=summary_result['key_points'],
            word_count=summary_result['word_count'],
            confidence=summary_result['confidence']
        )
        
        message = f"""✅ **Belge Özeti Oluşturuldu!**

━━━━━━━━━━━━━━━━━━━━━
{summary_result['summary']}

📊 **İstatistikler:**
• Kelime sayısı: {summary_result['word_count']}
• Özet güveni: %{summary_result['confidence']}

🔑 **Anahtar noktalar:**
"""
        for point in summary_result['key_points']:
            message += f"• {point}\n"
        
        message += f"\n📌 **1 hak tüketildi.**"
        
        # Kalan hakkı göster
        remaining = get_user_rights_direct(user_id)
        message += f"\n🔁 **Kalan hak:** {remaining}"
        
        keyboard = [[InlineKeyboardButton("◀️ İşlemlere Dön", callback_data="merhaba")]]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        await query.edit_message_text(message, reply_markup=reply_markup, parse_mode='Markdown')
        
    except Exception as e:
        logger.error(f"❌ Özetleme hatası: {e}")
        await query.edit_message_text(f"❌ Hata oluştu: {str(e)[:100]}")

async def handle_smart_validate(query, context, user_id, file_path, file_type, file_name, file_size):
    """Hata ve eksik kontrolü işlemi"""
    await query.edit_message_text("✅ **Hata ve eksik kontrolü yapılıyor...**\n\n⏳ Lütfen bekleyin...")
    
    try:
        # Metin çıkar
        text_content = extract_text_from_file(file_path, file_type)
        
        # Basit extracted_info oluştur
        extracted_info = {}
        try:
            from converters import detect_important_fields
            extracted_info = detect_important_fields(text_content)
        except:
            pass
        
        # Belge türünü belirle
        doc_type = file_type.lower()
        
        # Doğrula
        validation_result = validator.validate_document(text_content, extracted_info, doc_type)
        
        # Hak tüket
        rights_success = db.increase_validation_count(user_id)
        
        if not rights_success:
            await query.edit_message_text("❌ Doğrulama için yeterli hakkınız bulunmamaktadır.")
            return
        
        # Kaydet
        db.save_validation_record(
            user_id=user_id,
            file_name=file_name,
            is_valid=validation_result['is_valid'],
            issues=[i['message'] for i in validation_result['issues']],
            warnings_count=validation_result['warnings_count'],
            errors_count=validation_result['errors_count'],
            critical_count=validation_result['critical_count'],
            score=validation_result['score']
        )
        
        # Rapor oluştur
        report = validator.get_validation_report(validation_result)
        
        message = f"✅ **Doğrulama Raporu**\n\n{report}\n\n📌 **1 hak tüketildi.**"
        
        # Kalan hakkı göster
        remaining = get_user_rights_direct(user_id)
        message += f"\n🔁 **Kalan hak:** {remaining}"
        
        keyboard = [[InlineKeyboardButton("◀️ İşlemlere Dön", callback_data="merhaba")]]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        await query.edit_message_text(message, reply_markup=reply_markup, parse_mode='Markdown')
        
    except Exception as e:
        logger.error(f"❌ Doğrulama hatası: {e}")
        await query.edit_message_text(f"❌ Hata oluştu: {str(e)[:100]}")

async def handle_smart_all(query, context, user_id, file_path, file_type, file_name, file_size):
    """Tüm akıllı işlemleri tek seferde yap"""
    await query.edit_message_text(
        "🔁 **Tüm akıllı işlemler başlatılıyor...**\n\n"
        "📛 Akıllı isimlendirme (1 hak)\n"
        "📄 Belge türü tanıma (1 hak)\n"
        "✨ Akıllı dönüşüm (2 hak)\n"
        "📋 Belge özetleme (1 hak)\n"
        "✅ Hata kontrolü (1 hak)\n\n"
        "⏳ Toplam **5 hak** tükenecektir. Lütfen bekleyin..."
    )
    
    try:
        # Tüm işlemleri yap
        success, message, results = await converters.smart_process_file(
            input_path=file_path,
            output_path=None,  # Özel durum
            source_type=file_type,
            target_type=None,  # Özel durum
            user_id=user_id,
            db_instance=db
        )
        
        # Kalan hakkı göster
        remaining = get_user_rights_direct(user_id)
        
        final_message = f"✅ **Tüm işlemler tamamlandı!**\n\n{message}\n\n🔁 **Kalan hak:** {remaining}"
        
        keyboard = [[InlineKeyboardButton("◀️ İşlemlere Dön", callback_data="merhaba")]]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        await query.edit_message_text(final_message, reply_markup=reply_markup, parse_mode='Markdown')
        
    except Exception as e:
        logger.error(f"❌ Toplu işlem hatası: {e}")
        await query.edit_message_text(f"❌ Hata oluştu: {str(e)[:100]}")

async def action_cancel_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """İşlem iptal et"""
    query = update.callback_query
    await query.answer()
    
    keyboard = [[InlineKeyboardButton("👋 Merhaba", callback_data="merhaba")]]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await query.edit_message_text(
        "❌ **İşlem iptal edildi.**\n\n"
        "Ana menüye döndünüz.",
        reply_markup=reply_markup,
        parse_mode='Markdown'
    )

async def show_conversion_options(update: Update, context: ContextTypes.DEFAULT_TYPE, file_type):
    """Dosya türüne göre buton menüsü göster (klasik dönüşüm)"""
    
    options = CONVERSION_MAP.get(file_type, ['PDF'])
    
    keyboard = []
    for opt in options:
        display_name = DISPLAY_NAMES.get(opt, opt)
        callback_data = f"convert|{opt}"
        keyboard.append([InlineKeyboardButton(display_name, callback_data=callback_data)])
    
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await update.message.reply_text(
        f"📄 **Dosya Algılandı**\n\n"
        f"📁 Dosya türü: **{DISPLAY_NAMES.get(file_type, file_type)}**\n"
        f"🔄 Dönüştürülebilecek formatlar:\n\n"
        f"Lütfen dönüştürmek istediğiniz formatı seçin:",
        reply_markup=reply_markup,
        parse_mode='Markdown'
    )

async def convert_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Dönüştürme butonuna basıldığında - KLASİK VERSİYON"""
    query = update.callback_query
    await query.answer()
    
    if not query.data.startswith('convert'):
        return
    
    target = query.data.split('|')[1]
    
    # Dosya bilgilerini al
    file_path = context.user_data.get('current_file')
    file_name = context.user_data.get('file_name')
    file_type = context.user_data.get('file_type')
    file_size = context.user_data.get('file_size', 0)
    user_id = update.effective_user.id
    
    if not file_path or not os.path.exists(file_path):
        await query.edit_message_text("❌ **Dosya bulunamadı!**\nLütfen dosyayı tekrar yükleyin.")
        return
    
    # Çıktı dosyası adını oluştur
    base_name = os.path.splitext(file_name)[0]
    output_ext = EXTENSION_MAP.get(target, '.pdf')
    safe_name = utils.safe_filename(f"{user_id}_{base_name}_converted{output_ext}")
    output_path = f"temp/{safe_name}"
    
    # Dönüşüm başlangıç zamanı
    start_time = datetime.datetime.now()
    
    # Bekleme mesajı
    await query.edit_message_text(
        f"⏳ **Dosya dönüştürülüyor...**\n\n"
        f"📁 Kaynak: `{file_name}`\n"
        f"🔄 Hedef: **{DISPLAY_NAMES.get(target, target)}**\n\n"
        f"Bu işlem birkaç saniye sürebilir, lütfen bekleyin..."
    )
    
    # Dönüştürme işlemini yap (klasik)
    success, error = await converters.convert_file(file_path, output_path, file_type, target)
    
    # İşlem süresini hesapla
    end_time = datetime.datetime.now()
    processing_time = (end_time - start_time).total_seconds()
    
    if success and os.path.exists(output_path):
        # Hakkı azalt
        db.decrease_rights(user_id, 'direct')
        
        # Dönüşüm kaydını ekle
        db.save_conversion_record(
            user_id=user_id,
            file_name=file_name,
            file_size=file_size,
            source_format=file_type,
            target_format=target,
            conversion_type='direct',
            status='success',
            processing_time=processing_time,
            quality_score=95
        )
        
        # Kullanıcıya dönüştürülmüş dosyayı gönder
        with open(output_path, 'rb') as f:
            await query.message.reply_document(
                document=f,
                filename=f"{base_name}_converted{output_ext}",
                caption=f"✅ **Dönüştürme tamamlandı!**"
            )
        
        # Dönüşüm mesajını gönder
        await query.message.reply_text("✅ Dosya doğrudan dönüştürüldü. Yapı ve içerik korunmuştur.", parse_mode='Markdown')
        
        # Yeni hak miktarını al
        new_remaining = get_user_rights_direct(user_id)
        
        # Kalan hak bilgisi
        await query.message.reply_text(f"🔁 **Kalan hak:** {new_remaining}", parse_mode='Markdown')
        
        logger.info(f"✅ Klasik dönüşüm başarılı: {user_id} - {file_type} -> {target}")
        
        # Geçici dosyaları temizle
        utils.clean_temp_files(user_id, file_path, output_path)
        
    else:
        # Başarısız dönüşüm
        db.increase_failed_count(user_id)
        
        db.save_conversion_record(
            user_id=user_id,
            file_name=file_name,
            file_size=file_size,
            source_format=file_type,
            target_format=target,
            conversion_type='direct',
            status='failed',
            processing_time=processing_time,
            error_message=error
        )
        
        error_msg = error if error else "Bilinmeyen hata"
        await query.message.reply_text(
            f"❌ **Dönüştürme başarısız!**\n\n"
            f"⚠️ Hata: `{error_msg[:200]}`\n\n"
            f"📂 **Yeni dosyanızı bekliyorum...**"
        )
        
        logger.warning(f"⚠️ Klasik dönüşüm başarısız: {user_id} - {file_type} -> {target}")
        
        # Geçici dosyayı temizle
        utils.clean_temp_files(user_id, file_path)

# ========== ADMIN KOMUTLARI ==========
async def admin_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Admin komutu"""
    user_id = update.effective_user.id
    
    if user_id != ADMIN_ID:
        await update.message.reply_text("❌ Bu komutu kullanma yetkiniz yok.")
        logger.warning(f"⚠️ Yetkisiz admin erişimi: {user_id}")
        return
    
    # Bekleyen ödeme sayısını göster
    conn = sqlite3.connect('database/bot.db')
    cursor = conn.cursor()
    cursor.execute("SELECT COUNT(*) FROM pending_payments WHERE status = 'pending'")
    pending_count = cursor.fetchone()[0]
    conn.close()
    
    keyboard = [
        [InlineKeyboardButton("📊 Sistem Durumu", callback_data="admin_durum")],
        [InlineKeyboardButton("👥 Toplam Kullanıcı", callback_data="admin_kullanici")],
        [InlineKeyboardButton("📈 Bugünkü Dönüşümler", callback_data="admin_bugun")],
        [InlineKeyboardButton("✅ Başarılı Dönüşümler", callback_data="admin_basarili")],
        [InlineKeyboardButton("❌ Başarısız Dönüşümler", callback_data="admin_basarisiz")],
        [InlineKeyboardButton("📊 Analiz İstatistikleri", callback_data="admin_analysis")],
        [InlineKeyboardButton("📛 İsimlendirme İstatistikleri", callback_data="admin_naming")],
        [InlineKeyboardButton("📋 Özetleme İstatistikleri", callback_data="admin_summary")],
        [InlineKeyboardButton("✅ Doğrulama İstatistikleri", callback_data="admin_validation")],
        [InlineKeyboardButton("⭐ Kalite Optimizasyonu", callback_data="admin_quality")],
        [InlineKeyboardButton(f"💰 Bekleyen Ödemeler ({pending_count})", callback_data="admin_pending_payments")],
        [InlineKeyboardButton("🔍 Kullanıcı Sorgula", callback_data="admin_sorgula")]
    ]
    
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await update.message.reply_text(
        "👑 **Admin Paneli**\n\n"
        "Lütfen bir işlem seçin:",
        reply_markup=reply_markup,
        parse_mode='Markdown'
    )
    
    logger.info(f"👑 Admin paneli açıldı: {user_id}")

async def admin_button_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Admin butonlarını yönet"""
    query = update.callback_query
    await query.answer()
    
    user_id = update.effective_user.id
    if user_id != ADMIN_ID:
        await query.message.reply_text("❌ Yetkisiz erişim!")
        return
    
    stats = db.get_admin_stats()
    
    if query.data == "admin_durum" and stats:
        await query.message.reply_text(
            f"📊 **Sistem Durumu**\n\n"
            f"👥 Toplam Kullanıcı: `{stats['total_users']}`\n"
            f"📈 Bugünkü Dönüşüm: `{stats['today_conversions']}`\n"
            f"✅ Başarılı: `{stats['success_total']}`\n"
            f"❌ Başarısız: `{stats['failed_total']}`\n"
            f"📁 Toplam Dönüşüm: `{stats['success_total'] + stats['failed_total']}`\n"
            f"📊 Toplam Başarılı: `{stats['total_success']}`\n"
            f"📊 Toplam Başarısız: `{stats['total_failed']}`\n"
            f"📊 Aktif Kullanıcılar: `{stats.get('active_users', 0)}`",
            parse_mode='Markdown'
        )
    
    elif query.data == "admin_kullanici" and stats:
        await query.message.reply_text(f"👥 **Toplam Kayıtlı Kullanıcı:** `{stats['total_users']}`", parse_mode='Markdown')
    
    elif query.data == "admin_bugun" and stats:
        await query.message.reply_text(f"📈 **Bugünkü Dönüşümler:** `{stats['today_conversions']}`", parse_mode='Markdown')
    
    elif query.data == "admin_basarili" and stats:
        await query.message.reply_text(f"✅ **Başarılı Dönüşümler:** `{stats['success_total']}`", parse_mode='Markdown')
    
    elif query.data == "admin_basarisiz" and stats:
        await query.message.reply_text(f"❌ **Başarısız Dönüşümler:** `{stats['failed_total']}`", parse_mode='Markdown')
    
    elif query.data == "admin_analysis" and stats:
        await query.message.reply_text(
            f"📊 **Analiz İstatistikleri**\n\n"
            f"📈 Bugünkü Analiz: `{stats.get('today_analysis', 0)}`\n"
            f"📊 Toplam Analiz: `{stats.get('total_analysis', 0)}`\n"
            f"✨ Akıllı Düzenleme: `{stats.get('smart_edit_total', 0)}`\n\n"
            f"📋 **Haftalık İstatistikler:**\n{stats.get('weekly_stats', 'Veri yok')}",
            parse_mode='Markdown'
        )
    
    elif query.data == "admin_naming" and stats:
        await query.message.reply_text(
            f"📛 **İsimlendirme İstatistikleri**\n\n"
            f"📈 Bugünkü İsimlendirme: `{stats.get('today_naming', 0)}`\n"
            f"📊 Toplam İsimlendirme: `{stats.get('total_naming', 0)}`",
            parse_mode='Markdown'
        )
    
    elif query.data == "admin_summary" and stats:
        await query.message.reply_text(
            f"📋 **Özetleme İstatistikleri**\n\n"
            f"📈 Bugünkü Özetleme: `{stats.get('today_summaries', 0)}`\n"
            f"📊 Toplam Özetleme: `{stats.get('total_summaries', 0)}`",
            parse_mode='Markdown'
        )
    
    elif query.data == "admin_validation" and stats:
        await query.message.reply_text(
            f"✅ **Doğrulama İstatistikleri**\n\n"
            f"📈 Bugünkü Doğrulama: `{stats.get('today_validations', 0)}`\n"
            f"📊 Toplam Doğrulama: `{stats.get('total_validations', 0)}`",
            parse_mode='Markdown'
        )
    
    elif query.data == "admin_quality" and stats:
        await query.message.reply_text(
            f"⭐ **Kalite Optimizasyonu İstatistikleri**\n\n"
            f"📊 Bu özellik yeni eklendi. Yakında istatistikler eklenecek.",
            parse_mode='Markdown'
        )
    
    elif query.data == "admin_pending_payments":
        # Bekleyen ödemeleri göster
        conn = sqlite3.connect('database/bot.db')
        cursor = conn.cursor()
        cursor.execute('''
            SELECT id, username, package_name, amount, requested_at 
            FROM pending_payments 
            WHERE status = 'pending'
            ORDER BY requested_at DESC
        ''')
        pending = cursor.fetchall()
        conn.close()
        
        if pending:
            text = "💰 **BEKLEYEN ÖDEMELER**\n\n"
            for p in pending:
                text += f"• `#{p[0]}` - @{p[1]} - {p[2]} - {p[3]} TL - {p[4][:16]}\n"
            await query.message.reply_text(text, parse_mode='Markdown')
        else:
            await query.message.reply_text("✅ **Bekleyen ödeme yok.**", parse_mode='Markdown')
    
    elif query.data == "admin_sorgula":
        await query.message.reply_text(
            "🔍 **Kullanıcı Sorgulama**\n\n"
            "Sorgulamak istediğiniz kullanıcının Telegram ID'sini gönderin.\n\n"
            "Örnek: `123456789`"
        )
        context.user_data['awaiting_user_id'] = True

async def handle_user_query(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Kullanıcı ID sorgulamasını yönet"""
    if context.user_data.get('awaiting_user_id'):
        user_id = update.message.text.strip()
        
        try:
            user_id = int(user_id)
            conn = sqlite3.connect('database/bot.db')
            cursor = conn.cursor()
            
            cursor.execute('''
                SELECT user_id, username, first_name, package_type, remaining_rights, 
                       successful_conversions, failed_conversions, total_analysis,
                       total_smart_edits, total_naming, total_classification,
                       total_summaries, total_validations, registered_at
                FROM users WHERE user_id = ?
            ''', (user_id,))
            user = cursor.fetchone()
            
            if user:
                await update.message.reply_text(
                    f"👤 **Kullanıcı Bilgileri**\n\n"
                    f"🆔 ID: `{user[0]}`\n"
                    f"👤 Kullanıcı Adı: @{user[1] if user[1] else 'Yok'}\n"
                    f"📝 İsim: {user[2]}\n"
                    f"📦 Paket: {user[3]} Dosya\n"
                    f"🔁 Kalan Hak: {user[4]}\n"
                    f"✅ Başarılı: {user[5]}\n"
                    f"❌ Başarısız: {user[6]}\n"
                    f"📊 Analiz Sayısı: {user[7] or 0}\n"
                    f"✨ Akıllı Düzenleme: {user[8] or 0}\n"
                    f"📛 İsimlendirme: {user[9] or 0}\n"
                    f"📄 Sınıflandırma: {user[10] or 0}\n"
                    f"📋 Özetleme: {user[11] or 0}\n"
                    f"✅ Doğrulama: {user[12] or 0}\n"
                    f"📅 Kayıt: {user[13]}\n",
                    parse_mode='Markdown'
                )
                logger.info(f"🔍 Kullanıcı sorgulandı: {user_id}")
            else:
                await update.message.reply_text("❌ Kullanıcı bulunamadı.")
            
            conn.close()
            
        except ValueError:
            await update.message.reply_text("❌ Geçersiz ID formatı. Lütfen sadece rakam girin.")
        except Exception as e:
            logger.error(f"❌ Kullanıcı sorgulama hatası: {e}")
            await update.message.reply_text(f"❌ Hata: {str(e)}")
        
        context.user_data['awaiting_user_id'] = False

# ========== ANA FONKSİYON ==========
def main():
    """Botu başlat"""
    check_single_instance()  # Tek instance kontrolü
    
    print("🚀 Dosya Asistanı Bot başlatılıyor...")
    print("=" * 60)
    
    print(f"🔑 Token: {BOT_TOKEN[:15]}...")
    print(f"👑 Admin ID: {ADMIN_ID}")
    print(f"📁 Modüler yapı: 10 dosya aktif")
    print(f"🔄 Dönüşüm: converters.py (SÜPER GELİŞMİŞ)")
    print(f"🤖 Analiz: analyzer.py (AKTİF)")
    print(f"📛 İsimlendirme: naming.py (AKTİF)")
    print(f"📄 Sınıflandırma: classifier.py (AKTİF)")
    print(f"📋 Özetleme: summarizer.py (AKTİF)")
    print(f"✅ Doğrulama: validator.py (AKTİF)")
    print(f"⭐ Kalite: quality_optimizer.py (AKTİF)")
    print(f"💰 Ödeme: payments.py (TELEFONSUZ - KULLANICI ADI İLE)")
    print(f"📊 Loglama: bot.log, payments.log, database.log, converters.log, quality.log")
    print(f"⚡ 7/24 HİZMETİNİZDEYİZ!")
    
    # Veritabanını oluştur
    try:
        db.init_database()
        init_payments_table()
        print("✅ Veritabanı başarıyla oluşturuldu/güncellendi")
    except Exception as e:
        print(f"❌ Veritabanı hatası: {e}")
        # Acil durum çözümü
        conn = sqlite3.connect('database/bot.db')
        cursor = conn.cursor()
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS users (
                user_id INTEGER PRIMARY KEY,
                username TEXT,
                first_name TEXT,
                package_type TEXT DEFAULT '30',
                remaining_rights INTEGER DEFAULT 30,
                total_conversions INTEGER DEFAULT 0,
                successful_conversions INTEGER DEFAULT 0,
                failed_conversions INTEGER DEFAULT 0,
                last_package_date TEXT,
                registered_at TEXT
            )
        ''')
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS conversions (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER,
                file_name TEXT,
                file_size INTEGER,
                source_format TEXT,
                target_format TEXT,
                status TEXT,
                processing_time REAL,
                error_message TEXT,
                converted_at TEXT
            )
        ''')
        conn.commit()
        conn.close()
        print("✅ Veritabanı acil durumda oluşturuldu.")
    
    # Bot uygulamasını oluştur
    application = Application.builder().token(BOT_TOKEN).build()
    
    # Komut handler'ları
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("admin", admin_command))
    
    # Buton handler'ları - smart_menu kaldırıldı
    application.add_handler(CallbackQueryHandler(button_handler, pattern="^(merhaba|dosya_yukle|check_rights)$"))
    application.add_handler(CallbackQueryHandler(convert_handler, pattern="^convert"))
    application.add_handler(CallbackQueryHandler(smart_action_handler, pattern="^action_"))
    application.add_handler(CallbackQueryHandler(smart_convert_to_handler, pattern="^smart_convert_to"))
    application.add_handler(CallbackQueryHandler(action_cancel_handler, pattern="^action_cancel$"))
    application.add_handler(CallbackQueryHandler(admin_button_handler, pattern="^admin_"))
    
    # Akıllı işlem handler'ları (ana menüden çağrılmayacak ama dosya sonrası için kalmalı)
    application.add_handler(CallbackQueryHandler(smart_naming_handler, pattern="^smart_naming$"))
    application.add_handler(CallbackQueryHandler(smart_classify_handler, pattern="^smart_classify$"))
    application.add_handler(CallbackQueryHandler(smart_convert_handler, pattern="^smart_convert$"))
    application.add_handler(CallbackQueryHandler(smart_summarize_handler, pattern="^smart_summarize$"))
    application.add_handler(CallbackQueryHandler(smart_validate_handler, pattern="^smart_validate$"))
    application.add_handler(CallbackQueryHandler(smart_all_handler, pattern="^smart_all$"))
    application.add_handler(CallbackQueryHandler(smart_quality_handler, pattern="^smart_quality$"))
    
    # Kalite seçim handler'ları
    application.add_handler(CallbackQueryHandler(quality_handler, pattern="^quality_"))
    application.add_handler(CallbackQueryHandler(back_to_convert_handler, pattern="^back_to_convert$"))
    
    # Ödeme handler'ları
    application.add_handler(CallbackQueryHandler(show_packages, pattern="^show_packages$"))
    application.add_handler(CallbackQueryHandler(show_package_detail, pattern="^package_"))
    application.add_handler(CallbackQueryHandler(start_payment, pattern="^buy_"))
    application.add_handler(CallbackQueryHandler(confirm_payment, pattern="^confirm_payment_"))
    application.add_handler(CallbackQueryHandler(approve_payment, pattern="^approve_payment_"))
    application.add_handler(CallbackQueryHandler(reject_payment, pattern="^reject_payment_"))
    application.add_handler(CallbackQueryHandler(cancel_payment, pattern="^cancel_payment$"))
    application.add_handler(CallbackQueryHandler(back_to_main, pattern="^back_to_main$"))
    
    # Mesaj handler'ları
    application.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_user_query))
    
    print("✅ Bot yapılandırması tamamlandı.")
    print("=" * 60)
    print("🤖 Bot çalışıyor...")
    print("📱 Telegram: @dosya_asistani_bot")
    print("⚡ 7/24 HİZMETİNİZDEYİZ!")
    print("✨ Akıllı İşlemler: Dosya yüklendikten sonra gösterilir")
    print("⭐ Kalite Seviyeleri: Taslak/Standart/Profesyonel/Premium")
    print("🛑 Durdurmak: CTRL+C")
    print("=" * 60)
    
    try:
        application.run_polling()
    except Exception as e:
        print(f"❌ Bot çalışırken hata: {e}")
        logger.error(f"❌ Bot hatası: {e}")
    finally:
        print("👋 Bot durduruldu.")
        logger.info("👋 Bot durduruldu.")
        # Lock dosyasını temizle
        try:
            if os.path.exists(LOCK_FILE):
                os.remove(LOCK_FILE)
        except:
            pass

if __name__ == "__main__":
    main()
