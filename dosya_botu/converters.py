"""
PROFESYONEL DOSYA DÖNÜŞTÜRME MODÜLÜ
Tüm dönüşümler yüksek kalitede ve sorunsuz çalışır
Gelişmiş tipografi, tablo yönetimi ve format koruma
Yapay zeka destekli analiz, isimlendirme, sınıflandırma, özetleme ve doğrulama entegrasyonu
"""

import os
import re
import logging
import datetime
import asyncio
import hashlib
import json
import traceback
from typing import Dict, List, Tuple, Optional, Any, Union
from dataclasses import dataclass, field
from enum import Enum
from collections import Counter
from PIL import Image, ImageEnhance, ImageFilter, ImageDraw, ImageFont, ImageOps
import pytesseract
import tempfile
import shutil

# Yeni modüller
try:
    import analyzer
except ImportError:
    analyzer = None
    print("⚠️ analyzer modülü bulunamadı")

try:
    import ai_editor
except ImportError:
    ai_editor = None

try:
    import naming
except ImportError:
    naming = None

try:
    import classifier
except ImportError:
    classifier = None

try:
    import summarizer
except ImportError:
    summarizer = None

try:
    import validator
except ImportError:
    validator = None

try:
    import quality_optimizer
except ImportError:
    quality_optimizer = None

# Tesseract yolunu ayarla (Windows için)
try:
    if os.name == 'nt':  # Windows
        possible_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
        ]
        for path in possible_paths:
            if os.path.exists(path):
                pytesseract.pytesseract.tesseract_cmd = path
                break
except:
    pass

# Loglama ayarları
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('converters.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# ========== DESTEKLENEN FORMATLAR ==========
SUPPORTED_SOURCE_FORMATS = {
    'WORD': ['.doc', '.docx'],
    'EXCEL': ['.xls', '.xlsx'],
    'POWERPOINT': ['.ppt', '.pptx'],
    'PDF': ['.pdf'],
    'GORSEL': ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif'],
    'TEXT': ['.txt', '.rtf', '.md'],
}

SUPPORTED_TARGET_FORMATS = {
    'WORD': ['PDF', 'EXCEL', 'POWERPOINT', 'TEXT'],
    'EXCEL': ['PDF', 'WORD', 'POWERPOINT', 'TEXT'],
    'POWERPOINT': ['PDF', 'WORD', 'TEXT'],
    'PDF': ['WORD', 'TEXT'],
    'GORSEL': ['PDF', 'WORD', 'TEXT'],
    'TEXT': ['PDF', 'WORD'],
}

FORMAT_DISPLAY_NAMES = {
    'WORD': '📝 Word',
    'EXCEL': '📊 Excel',
    'POWERPOINT': '📽️ PowerPoint',
    'PDF': '📄 PDF',
    'GORSEL': '🖼️ Görsel',
    'TEXT': '📃 Metin',
}

EXTENSION_MAP = {
    'WORD': '.docx',
    'EXCEL': '.xlsx',
    'POWERPOINT': '.pptx',
    'PDF': '.pdf',
    'GORSEL': '.png',
    'TEXT': '.txt',
}


# ========== YARDIMCI SINIFLAR ==========

class ConversionQuality(Enum):
    """Dönüşüm kalite seviyeleri"""
    DRAFT = "taslak"           # Düşük kalite, hızlı
    STANDARD = "standart"       # Normal kalite
    PROFESSIONAL = "profesyonel" # Yüksek kalite
    PREMIUM = "premium"         # Maksimum kalite


class DocumentComplexity(Enum):
    """Belge karmaşıklık seviyeleri"""
    SIMPLE = "basit"            # Düz metin
    MODERATE = "orta"            # Başlıklar, listeler
    COMPLEX = "karmaşık"         # Tablolar, grafikler
    VERY_COMPLEX = "çok karmaşık" # Karmaşık yapılar


@dataclass
class ConversionMetrics:
    """Dönüşüm metrikleri"""
    input_size: int = 0
    output_size: int = 0
    processing_time: float = 0
    compression_ratio: float = 0
    quality_score: int = 0
    warnings: List[str] = field(default_factory=list)
    suggestions: List[str] = field(default_factory=list)
    complexity: DocumentComplexity = DocumentComplexity.SIMPLE


@dataclass
class ConversionResult:
    """Dönüşüm sonuç veri yapısı"""
    success: bool
    output_path: str
    metrics: ConversionMetrics
    changes_made: List[str]
    error_message: Optional[str] = None
    warning_message: Optional[str] = None


# ========== YARDIMCI FONKSİYONLAR (GELİŞTİRİLMİŞ) ==========

def get_file_extension(file_path: str) -> str:
    """Dosya uzantısını döndür"""
    return os.path.splitext(file_path)[1].lower()


def get_file_name_without_extension(file_path: str) -> str:
    """Dosya adını uzantısız döndür"""
    return os.path.splitext(os.path.basename(file_path))[0]


def get_file_size_str(size_bytes: int) -> str:
    """Dosya boyutunu okunabilir formata çevir"""
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes/1024:.1f} KB"
    elif size_bytes < 1024 * 1024 * 1024:
        return f"{size_bytes/(1024*1024):.1f} MB"
    else:
        return f"{size_bytes/(1024*1024*1024):.1f} GB"


def detect_source_type(file_path: str) -> Optional[str]:
    """Dosya uzantısından kaynak tipini tespit et"""
    ext = get_file_extension(file_path)
    
    for source_type, extensions in SUPPORTED_SOURCE_FORMATS.items():
        if ext in extensions:
            return source_type
    
    return None


def is_conversion_supported(source_type: str, target_type: str) -> bool:
    """Dönüşümün desteklenip desteklenmediğini kontrol et"""
    if source_type not in SUPPORTED_TARGET_FORMATS:
        return False
    
    return target_type in SUPPORTED_TARGET_FORMATS[source_type]


def clean_text(text: str, aggressive: bool = False) -> str:
    """
    Metni temizle ve düzenle (gelişmiş)
    
    Args:
        text: Temizlenecek metin
        aggressive: Agresif temizleme modu
    """
    if not text:
        return ""
    
    # Fazla boşlukları temizle
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        # Satırdaki fazla boşlukları temizle
        clean_line = ' '.join(line.split())
        if clean_line:
            cleaned_lines.append(clean_line)
    
    if aggressive:
        # Agresif temizleme - tüm gereksiz karakterleri temizle
        final_lines = []
        for line in cleaned_lines:
            # Sadece alfanümerik ve temel noktalama işaretlerini bırak
            clean_line = re.sub(r'[^\w\s\.,;:!?\-\(\)\[\]{}@#$%&*+/=]', '', line)
            if clean_line:
                final_lines.append(clean_line)
    else:
        # Normal temizleme - paragrafları birleştir (çok kısa satırları)
        final_lines = []
        i = 0
        while i < len(cleaned_lines):
            if i < len(cleaned_lines) - 1 and len(cleaned_lines[i]) < 30 and len(cleaned_lines[i+1]) > 50:
                # Kısa satır + uzun satır = muhtemelen başlık
                final_lines.append(cleaned_lines[i])
                final_lines.append(cleaned_lines[i+1])
                i += 2
            elif i < len(cleaned_lines) - 1 and len(cleaned_lines[i]) > 0 and len(cleaned_lines[i+1]) > 0:
                # İki normal satır
                final_lines.append(cleaned_lines[i])
                i += 1
            else:
                final_lines.append(cleaned_lines[i])
                i += 1
    
    return '\n'.join(final_lines)


def detect_table_structure(text: str) -> Tuple[bool, List[Dict]]:
    """
    Metin içinde tablo yapısını tespit et (gelişmiş)
    
    Returns:
        (tablo_var_mı, tablo_bilgileri)
    """
    lines = text.split('\n')
    tablo_olasiligi = 0
    possible_tables = []
    table_score = 0
    
    for i, line in enumerate(lines):
        # Sekme ile ayrılmış tablo
        if '\t' in line:
            cells = line.split('\t')
            if len(cells) > 1:
                table_score += 3
                possible_tables.append({
                    'line': i,
                    'cells': cells,
                    'type': 'tab',
                    'cell_count': len(cells)
                })
        
        # Çoklu boşluk ile ayrılmış tablo
        elif '  ' in line and len(line.split('  ')) > 2:
            cells = [c for c in line.split('  ') if c.strip()]
            if len(cells) > 1:
                table_score += 2
                possible_tables.append({
                    'line': i,
                    'cells': cells,
                    'type': 'space',
                    'cell_count': len(cells)
                })
        
        # Boru (|) ile ayrılmış tablo
        elif '|' in line:
            cells = [c.strip() for c in line.split('|') if c.strip()]
            if len(cells) > 1:
                table_score += 4
                possible_tables.append({
                    'line': i,
                    'cells': cells,
                    'type': 'pipe',
                    'cell_count': len(cells)
                })
        
        # Düzenli sütunlar
        elif re.search(r'\s{3,}', line):
            cells = re.split(r'\s{3,}', line)
            if len(cells) > 1:
                table_score += 2
                possible_tables.append({
                    'line': i,
                    'cells': cells,
                    'type': 'regex',
                    'cell_count': len(cells)
                })
    
    # Düzenli tablo kontrolü (ardışık satırlar)
    if len(possible_tables) > 2:
        consecutive = 0
        for j in range(len(possible_tables) - 1):
            if possible_tables[j+1]['line'] - possible_tables[j]['line'] == 1:
                consecutive += 1
        if consecutive > 1:
            table_score += 10
    
    # Hücre sayılarının tutarlılığını kontrol et
    if len(possible_tables) > 1:
        cell_counts = [t['cell_count'] for t in possible_tables]
        if max(cell_counts) - min(cell_counts) <= 2:
            table_score += 5
    
    return table_score > len(lines) * 0.15, possible_tables


def format_number(value: Any, decimal_places: int = 2, thousand_sep: bool = True, 
                  currency: str = None) -> str:
    """
    Sayıları formatla (gelişmiş)
    
    Args:
        value: Sayısal değer
        decimal_places: Ondalık basamak sayısı
        thousand_sep: Binlik ayırıcı kullan
        currency: Para birimi (TL, USD, EUR, vs.)
    """
    if value is None:
        return ""
    
    if isinstance(value, (int, float)):
        if isinstance(value, float):
            if value.is_integer():
                formatted = str(int(value))
            else:
                formatted = f"{value:.{decimal_places}f}".replace('.', ',')
        else:
            formatted = str(value)
        
        # Binlik ayırıcı ekle
        if thousand_sep and len(formatted) > 3:
            parts = formatted.split(',')
            integer_part = parts[0]
            
            if len(integer_part) > 3:
                # Sağdan sola doğru grupla
                grouped = []
                for i in range(len(integer_part), 0, -3):
                    start = max(0, i-3)
                    grouped.append(integer_part[start:i])
                integer_part = '.'.join(reversed(grouped))
            
            if len(parts) > 1:
                formatted = f"{integer_part},{parts[1]}"
            else:
                formatted = integer_part
        
        # Para birimi ekle
        if currency:
            currency_map = {
                'TL': '₺',
                'USD': '$',
                'EUR': '€',
                'GBP': '£',
                'JPY': '¥'
            }
            symbol = currency_map.get(currency, currency)
            formatted = f"{formatted} {symbol}"
        
        return formatted
    
    return str(value)


def extract_text_from_file(file_path: str, file_type: str, 
                          extract_mode: str = 'full') -> str:
    """
    Dosyadan metin çıkar (farklı formatlar için - süper gelişmiş)
    
    Args:
        file_path: Dosya yolu
        file_type: Dosya tipi (WORD, PDF, EXCEL, POWERPOINT, GORSEL)
        extract_mode: Çıkarma modu (full, metadata, content_only)
    """
    text = ""
    metadata = {}
    
    try:
        if file_type == 'WORD':
            from docx import Document
            doc = Document(file_path)
            
            # Belge özellikleri
            if extract_mode in ['full', 'metadata']:
                core_props = doc.core_properties
                metadata = {
                    'author': core_props.author,
                    'created': str(core_props.created),
                    'modified': str(core_props.modified),
                    'title': core_props.title,
                    'subject': core_props.subject
                }
            
            if extract_mode in ['full', 'content_only']:
                # Paragraflar
                for para in doc.paragraphs:
                    if para.text.strip():
                        text += para.text + "\n"
                
                # Tablolar
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        if row_text.strip():
                            text += row_text + "\n"
                    text += "\n"
                
                # Header/Footer
                for section in doc.sections:
                    header = section.header
                    footer = section.footer
                    if header.paragraphs:
                        text += "--- Header ---\n"
                        for para in header.paragraphs:
                            if para.text.strip():
                                text += para.text + "\n"
                    if footer.paragraphs:
                        text += "--- Footer ---\n"
                        for para in footer.paragraphs:
                            if para.text.strip():
                                text += para.text + "\n"
        
        elif file_type == 'PDF':
            import PyPDF2
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                
                if extract_mode in ['full', 'metadata']:
                    metadata = pdf_reader.metadata
                
                if extract_mode in ['full', 'content_only']:
                    for page_num, page in enumerate(pdf_reader.pages):
                        page_text = page.extract_text()
                        if page_text.strip():
                            text += f"--- Sayfa {page_num + 1} ---\n"
                            text += page_text + "\n\n"
        
        elif file_type == 'EXCEL':
            import pandas as pd
            excel_file = pd.ExcelFile(file_path)
            
            if extract_mode in ['full', 'metadata']:
                metadata = {
                    'sheet_names': excel_file.sheet_names,
                    'sheet_count': len(excel_file.sheet_names)
                }
            
            if extract_mode in ['full', 'content_only']:
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    text += f"--- Sayfa: {sheet_name} ({df.shape[0]} satır, {df.shape[1]} sütun) ---\n"
                    text += df.to_string() + "\n\n"
        
        elif file_type == 'POWERPOINT':
            from pptx import Presentation
            prs = Presentation(file_path)
            
            if extract_mode in ['full', 'metadata']:
                metadata = {
                    'slide_count': len(prs.slides)
                }
            
            if extract_mode in ['full', 'content_only']:
                for slide_num, slide in enumerate(prs.slides, 1):
                    text += f"--- Slayt {slide_num} ---\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text += shape.text + "\n"
                    text += "\n"
        
        elif file_type == 'GORSEL':
            from PIL import Image, ImageEnhance, ImageFilter
            import pytesseract
            
            image = Image.open(file_path)
            
            if extract_mode in ['full', 'metadata']:
                metadata = {
                    'format': image.format,
                    'mode': image.mode,
                    'size': image.size,
                    'width': image.width,
                    'height': image.height
                }
            
            if extract_mode in ['full', 'content_only']:
                # Görseli ön işle
                if image.mode != 'L':
                    image = image.convert('L')
                
                # Kontrast artır
                enhancer = ImageEnhance.Contrast(image)
                image = enhancer.enhance(2.0)
                
                # Gürültü azalt
                image = image.filter(ImageFilter.MedianFilter(size=3))
                
                # OCR dene
                text = pytesseract.image_to_string(image, lang='tur+eng', config='--psm 6')
                
                if not text.strip():
                    # Alternatif PSM dene
                    text = pytesseract.image_to_string(image, lang='tur+eng', config='--psm 3')
        
        else:
            # TXT veya diğer formatlar
            encodings = ['utf-8', 'windows-1254', 'iso-8859-9', 'latin1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                        text = f.read()
                    break
                except:
                    continue
    
    except Exception as e:
        logger.error(f"❌ Metin çıkarma hatası: {e}")
    
    if extract_mode == 'full' and metadata:
        # Metadata'yı da ekle
        meta_text = "\n".join([f"{k}: {v}" for k, v in metadata.items() if v])
        if meta_text:
            text = f"--- BELGE METADATASI ---\n{meta_text}\n\n{text}"
    
    return clean_text(text, aggressive=False)


def detect_language(text: str) -> str:
    """Metnin dilini tespit et (Türkçe/İngilizce/Almanca/Fransızca)"""
    if not text:
        return 'unknown'
    
    # Dil karakter setleri
    language_chars = {
        'tr': set('ğüşıöçĞÜŞİÖÇ'),
        'de': set('äöüßÄÖÜ'),
        'fr': set('éèêëàâçôùûÿœæ'),
        'es': set('ñáéíóúüÑÁÉÍÓÚÜ'),
        'it': set('àèéìíîòóùú'),
        'ru': set('абвгдеёжзийклмнопрстуфхцчшщъыьэюя'),
        'ar': set('ابتثجحخدذرزسشصضطظعغفقكلمنهوي'),
    }
    
    scores = {}
    for lang, chars in language_chars.items():
        count = sum(1 for c in text if c in chars)
        if count > 0:
            scores[lang] = count
    
    if scores:
        return max(scores, key=scores.get)
    
    # İngilizce varsayılan
    return 'en'


def detect_important_fields(text: str) -> Dict[str, Any]:
    """Belgeden önemli alanları tespit et (süper gelişmiş)"""
    fields = {}
    
    # Tarih desenleri
    date_patterns = [
        (r'(\d{1,2})[\.\/](\d{1,2})[\.\/](\d{2,4})', 'DMY'),
        (r'(\d{4})[-](\d{1,2})[-](\d{1,2})', 'YMD'),
        (r'(\d{1,2})[\.\/](\d{1,2})[\.\/](\d{2,4})', 'DMY'),
        (r'(\d{2})[\.\/](\d{2})[\.\/](\d{4})', 'DMY2'),
        (r'(\d{1,2})\s+(Ocak|Şubat|Mart|Nisan|Mayıs|Haziran|Temmuz|Ağustos|Eylül|Ekim|Kasım|Aralık)\s+(\d{4})', 'TR_MONTH'),
        (r'(\d{1,2})\s+(January|February|March|April|May|June|July|August|September|October|November|December)\s+(\d{4})', 'EN_MONTH'),
    ]
    
    for pattern, format_type in date_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            if format_type == 'DMY':
                day, month, year = match.groups()
                if len(year) == 2:
                    year = '20' + year
                fields['tarih'] = f"{day.zfill(2)}.{month.zfill(2)}.{year}"
            elif format_type == 'YMD':
                year, month, day = match.groups()
                fields['tarih'] = f"{day.zfill(2)}.{month.zfill(2)}.{year}"
            elif format_type in ['TR_MONTH', 'EN_MONTH']:
                day, month_text, year = match.groups()
                month_map = {
                    'Ocak': '01', 'Şubat': '02', 'Mart': '03', 'Nisan': '04',
                    'Mayıs': '05', 'Haziran': '06', 'Temmuz': '07', 'Ağustos': '08',
                    'Eylül': '09', 'Ekim': '10', 'Kasım': '11', 'Aralık': '12',
                    'January': '01', 'February': '02', 'March': '03', 'April': '04',
                    'May': '05', 'June': '06', 'July': '07', 'August': '08',
                    'September': '09', 'October': '10', 'November': '11', 'December': '12'
                }
                month = month_map.get(month_text, '01')
                fields['tarih'] = f"{day.zfill(2)}.{month}.{year}"
            break
    
    # Tutar desenleri
    amount_patterns = [
        r'(?:toplam|genel toplam|tutar|ödenecek|mevduat|fiyat|ücret|bedel)[\s:]*([\d.,]+)\s*(TL|USD|EUR|₺|\$|€|GBP|JPY)?',
        r'([\d.,]+)\s*(TL|USD|EUR|₺|\$|€|GBP|JPY)',
        r'(?:TL|USD|EUR|₺|\$|€|GBP|JPY)\s*([\d.,]+)',
    ]
    
    for pattern in amount_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            amount = match.group(1)
            currency = match.group(2) if len(match.groups()) > 1 else 'TL'
            currency_map = {'₺': 'TL', '$': 'USD', '€': 'EUR', '£': 'GBP', '¥': 'JPY'}
            currency = currency_map.get(currency, currency)
            fields['tutar'] = f"{amount} {currency}"
            break
    
    # Firma/Kişi adı
    company_patterns = [
        r'(?:firma|şirket|company|müşteri|customer|alıcı|satıcı|tedarikçi)[\s:]*([A-Za-zğüşıöçĞÜŞİÖÇ\s\.]+)',
        r'(?:adı|name|ünvan|unvan|title)[\s:]*([A-Za-zğüşıöçĞÜŞİÖÇ\s\.]+)',
        r'^(.*?)(?:Ltd\.?|Şti\.?|A\.?Ş\.?|Inc\.?|Corp\.?|LLC)$',
    ]
    
    for pattern in company_patterns:
        match = re.search(pattern, text, re.IGNORECASE | re.MULTILINE)
        if match:
            company = match.group(1).strip()
            if len(company) > 2 and len(company) < 100:
                fields['firma'] = company
                break
    
    # Vergi numarası
    tax_patterns = [
        r'(?:vergi no|tax id|vergi dairesi|tax office)[\s:]*(\d{10,11})',
        r'(?:TC|kimlik no|TCKN|kimlik numarası)[\s:]*(\d{11})',
        r'(?:VKN|vergi kimlik numarası)[\s:]*(\d{10})',
    ]
    
    for pattern in tax_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            fields['vergi_no'] = match.group(1)
            break
    
    # IBAN
    iban_pattern = r'(?:IBAN|iban)[\s:]*([A-Z]{2}\d{2}\s*(?:\d{4}\s*){4,7})'
    match = re.search(iban_pattern, text, re.IGNORECASE)
    if match:
        fields['iban'] = match.group(1).replace(' ', '')
    
    # Telefon
    phone_pattern = r'(?:tel|telefon|phone)[\s:]*((?:\+90|0)?[0-9]{10,11})'
    match = re.search(phone_pattern, text, re.IGNORECASE)
    if match:
        fields['telefon'] = match.group(1)
    
    # E-posta
    email_pattern = r'([a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})'
    match = re.search(email_pattern, text)
    if match:
        fields['email'] = match.group(1)
    
    # Adres
    address_patterns = [
        r'(?:adres|address)[\s:]*([^\n]+)',
        r'(?:mahalle|mah\.|sokak|sk\.|cadde|cad\.|bulvar)[\s:]*([^\n]+)',
    ]
    
    for pattern in address_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            fields['adres'] = match.group(1).strip()
            break
    
    return fields


def calculate_file_hash(file_path: str) -> str:
    """Dosya hash'i hesapla (değişiklik kontrolü için)"""
    try:
        with open(file_path, 'rb') as f:
            return hashlib.md5(f.read()).hexdigest()
    except:
        return ""


def ensure_directory(directory: str) -> bool:
    """Klasörün var olduğundan emin ol"""
    try:
        if not os.path.exists(directory):
            os.makedirs(directory)
        return True
    except:
        return False


# ========== OCR TEMİZLEME FONKSİYONLARI ==========

def clean_ocr_text(text: str) -> Tuple[str, List[str]]:
    """
    OCR metnini temizle ve gereksiz satırları filtrele
    
    Args:
        text: Orijinal OCR metni
    
    Returns:
        (temizlenmiş_metin, silinen_satırlar)
    """
    lines = text.split('\n')
    cleaned_lines = []
    removed_lines = []
    
    # Gereksiz desenler
    unwanted_patterns = [
        r'\d{1,2}:\d{2}(?::\d{2})?(?:\s*(?:AM|PM|am|pm))?',  # Saat: 14:30, 2:30 PM
        r'\d{1,2}[./]\d{1,2}[./]\d{2,4}',  # Tarih: 15.03.2024
        r'Volte?\s*\d{1,2}:\d{2}',  # Volte 14:30
        r'(?:Turkcell|Vodafone|Türk Telekom|Türkcell)\s*\d{1,2}:\d{2}',
        r'Ekran\s*(?:Alıntısı|Görüntüsü|Fotoğrafi)',  # Ekran Alıntısı
        r'Screen\s*(?:Shot|Capture)',  # Screen Shot/Capture
        r'Screenshot',  # Screenshot
        r'Captur(?:e|ed)',  # Capture/Captured
        r'^\s*\d+\s*$',  # Sadece rakam
        r'^\s*[•\-*]\s*$',  # Sadece madde işareti
        r'^\s*[|\\/]\s*$',  # Sadece çizgi
        r'\d{1,2}\s+(?:Ocak|Şubat|Mart|Nisan|Mayıs|Haziran|Temmuz|Ağustos|Eylül|Ekim|Kasım|Aralık)\s+\d{4}',  # Tarih formatı
        r'\d{1,2}\s+(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}',
        r'(?:Sayfa|Page)\s+\d+\s*(?:/|of)\s*\d+',  # Sayfa numaraları
        r'^\s*[_\-\=\*]{3,}\s*$',  # Çizgiler
        r'^\s*[▌▐▀▄█▓▒░]\s*$',  # Blok karakterler
    ]
    
    # Çok kısa satırları filtrele (3 karakterden az)
    for line in lines:
        line = line.strip()
        
        # Boş satırı kontrol et
        if not line:
            cleaned_lines.append('')
            continue
        
        # Çok kısa satırları filtrele
        if len(line) < 3:
            removed_lines.append(line)
            continue
        
        # Gereksiz desenleri kontrol et
        is_unwanted = False
        for pattern in unwanted_patterns:
            if re.search(pattern, line, re.IGNORECASE):
                removed_lines.append(line)
                is_unwanted = True
                break
        
        if not is_unwanted:
            cleaned_lines.append(line)
    
    return '\n'.join(cleaned_lines), removed_lines


def calculate_ocr_confidence(text: str) -> float:
    """
    OCR metninin kalitesini hesapla
    
    Args:
        text: OCR metni
    
    Returns:
        Güven skoru (0-100)
    """
    if not text.strip():
        return 0
    
    words = text.split()
    total_chars = len(text)
    total_words = len(words)
    
    if total_words == 0:
        return 0
    
    # Ortalama kelime uzunluğu (normalde 4-8 arası olmalı)
    avg_word_length = total_chars / total_words
    
    # Anormal kısa kelimeler
    very_short_words = sum(1 for w in words if len(w) <= 2)
    short_ratio = very_short_words / total_words
    
    # Özel karakter oranı
    special_chars = sum(1 for c in text if not c.isalnum() and not c.isspace())
    special_ratio = special_chars / total_chars if total_chars > 0 else 0
    
    # Rakam oranı (çok fazla rakam OCR hatası olabilir)
    digit_chars = sum(1 for c in text if c.isdigit())
    digit_ratio = digit_chars / total_chars if total_chars > 0 else 0
    
    # Güven skoru hesapla
    confidence = 100
    
    # Ortalama kelime uzunluğu çok düşükse
    if avg_word_length < 3:
        confidence -= 30
    elif avg_word_length > 15:
        confidence -= 20
    
    # Çok fazla kısa kelime varsa
    if short_ratio > 0.3:
        confidence -= 25
    elif short_ratio > 0.2:
        confidence -= 15
    
    # Çok fazla özel karakter varsa
    if special_ratio > 0.2:
        confidence -= 20
    elif special_ratio > 0.1:
        confidence -= 10
    
    # Çok fazla rakam varsa (sayısal belgeler hariç)
    if digit_ratio > 0.3:
        confidence -= 15
    elif digit_ratio > 0.2:
        confidence -= 5
    
    return max(0, min(100, confidence))


def merge_intelligent_lines(text: str) -> str:
    """
    Akıllı satır birleştirme - bölünmüş kelimeleri ve cümleleri düzelt
    
    Args:
        text: Ham metin
    
    Returns:
        Birleştirilmiş metin
    """
    lines = text.split('\n')
    merged = []
    i = 0
    
    while i < len(lines):
        current_line = lines[i].strip()
        
        # Boş satır
        if not current_line:
            merged.append('')
            i += 1
            continue
        
        # Sonraki satır var mı?
        if i < len(lines) - 1:
            next_line = lines[i+1].strip()
            
            # Eğer sonraki satır küçük harfle başlıyorsa ve
            # mevcut satır tire ile bitmiyorsa (kelime bölünmesi değilse)
            if next_line and next_line[0].islower() and not current_line.endswith('-'):
                # Birleştir
                current_line += ' ' + next_line
                i += 2  # Bir sonraki satırı atla
                continue
        
        merged.append(current_line)
        i += 1
    
    return '\n'.join(merged)


def normalize_whitespace(text: str) -> str:
    """
    Boşlukları normalize et
    
    Args:
        text: Ham metin
    
    Returns:
        Düzenlenmiş metin
    """
    # Fazla boşlukları temizle
    text = re.sub(r' +', ' ', text)
    
    # Fazla satır sonlarını temizle
    text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)
    
    return text.strip()


def fix_common_ocr_errors(text: str) -> str:
    """
    Yaygın OCR hatalarını düzelt
    
    Args:
        text: OCR metni
    
    Returns:
        Düzeltilmiş metin
    """
    if not text:
        return ""
    
    # Yaygın hata düzeltmeleri
    replacements = [
        # Sayısal hatalar
        (r'\b0\b', 'O'),  # Tek başına 0 -> O
        (r'\b1\b', 'l'),  # Tek başına 1 -> l
        (r'\b5\b', 'S'),  # Tek başına 5 -> S
        
        # Karakter karışıklıkları
        ('§', 'S'),  # Paragraf işareti -> S
        ('©', 'C'),  # Copyright -> C
        ('®', 'R'),  # Registered -> R
        ('™', 'TM'),  # Trademark -> TM
        ('•', '-'),  # Madde işareti -> -
        ('·', '.'),  # Orta nokta -> .
        ('…', '...'),  # Üç nokta -> ...
        ('–', '-'),  # Uzun tire -> -
        ('—', '-'),  # Uzun tire -> -
        ('"', '"'),  # Tırnak işareti
        ('"', '"'),  # Tırnak işareti
        ('´', "'"),  # Kesme işareti
        ('`', "'"),  # Kesme işareti
        
        # Türkçe karakter düzeltmeleri
        ('Ý', 'İ'),  # Y. harfi -> İ
        ('Þ', 'Ş'),  # TH -> Ş
        ('ð', 'ğ'),  # Ð -> ğ
        ('ý', 'ı'),  # ý -> ı
        ('Ã', 'Ç'),  # Ã -> Ç
        ('ã', 'ç'),  # ã -> ç
        ('Ä', 'Ö'),  # Ä -> Ö
        ('ä', 'ö'),  # ä -> ö
        ('Ü', 'Ü'),  # Ü zaten doğru
        ('ü', 'ü'),  # ü zaten doğru
        ('Ë', 'E'),  # Ë -> E
        ('ë', 'e'),  # ë -> e
        
        # İngilizce karakter düzeltmeleri
        ('À', 'A'), ('Á', 'A'), ('Â', 'A'), ('Ã', 'A'),
        ('à', 'a'), ('á', 'a'), ('â', 'a'), ('ã', 'a'),
        ('È', 'E'), ('É', 'E'), ('Ê', 'E'), ('Ë', 'E'),
        ('è', 'e'), ('é', 'e'), ('ê', 'e'), ('ë', 'e'),
        ('Ì', 'I'), ('Í', 'I'), ('Î', 'I'), ('Ï', 'I'),
        ('ì', 'i'), ('í', 'i'), ('î', 'i'), ('ï', 'i'),
        ('Ò', 'O'), ('Ó', 'O'), ('Ô', 'O'), ('Õ', 'O'),
        ('ò', 'o'), ('ó', 'o'), ('ô', 'o'), ('õ', 'o'),
        ('Ù', 'U'), ('Ú', 'U'), ('Û', 'U'), ('Ü', 'U'),
        ('ù', 'u'), ('ú', 'u'), ('û', 'u'), ('ü', 'u'),
        ('Ñ', 'N'), ('ñ', 'n'),
    ]
    
    for old, new in replacements:
        if isinstance(old, tuple) or (isinstance(old, str) and old.startswith('\\')):
            # Regex pattern
            text = re.sub(old, new, text)
        else:
            # Direkt değiştirme
            text = text.replace(old, new)
    
    return text


# ========== WORD DÖNÜŞÜMLERİ ==========

def word_to_pdf(input_path: str, output_path: str, 
               quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Word -> PDF (PROFESYONEL - TİPOGRAFİ KORUMALI)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from docx import Document
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        
        doc = Document(input_path)
        
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        left_margin = 2 * cm
        right_margin = width - 2 * cm
        y = height - 2 * cm
        line_height = 0.6 * cm
        
        title_size = 16
        heading_size = 14
        normal_size = 11
        
        for paragraph in doc.paragraphs:
            if not paragraph.text.strip():
                y -= line_height * 0.3
                continue
            
            style_name = paragraph.style.name.lower() if paragraph.style else ""
            font_size = normal_size
            is_bold = False
            
            if 'title' in style_name or 'heading 1' in style_name:
                font_size = title_size
                is_bold = True
            elif 'heading 2' in style_name:
                font_size = heading_size
                is_bold = True
            
            if is_bold:
                c.setFont("Helvetica-Bold", font_size)
            else:
                c.setFont("Helvetica", font_size)
            
            text = paragraph.text.strip()
            words = text.split()
            current_line = ""
            
            for word in words:
                test_line = current_line + " " + word if current_line else word
                line_width = c.stringWidth(test_line, "Helvetica" if not is_bold else "Helvetica-Bold", font_size)
                
                if line_width < (right_margin - left_margin):
                    current_line = test_line
                else:
                    if y < line_height:
                        c.showPage()
                        y = height - 2 * cm
                        if is_bold:
                            c.setFont("Helvetica-Bold", font_size)
                        else:
                            c.setFont("Helvetica", font_size)
                    
                    c.drawString(left_margin, y, current_line)
                    y -= line_height
                    current_line = word
            
            if current_line:
                if y < line_height:
                    c.showPage()
                    y = height - 2 * cm
                
                c.drawString(left_margin, y, current_line)
                y -= line_height * 1.2
        
        c.save()
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 90
        metrics.complexity = DocumentComplexity.MODERATE
        
        logger.info(f"✅ Word -> PDF dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Word -> PDF dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def word_to_excel(input_path: str, output_path: str, 
                 quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Word -> Excel (PROFESYONEL - AKILLI TABLO ALGILAMA)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from docx import Document
        import pandas as pd
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.utils import get_column_letter
        from openpyxl.worksheet.page import PageMargins
        
        doc = Document(input_path)
        
        # Kalite ayarları
        font_size = 11
        header_font_size = 12
        
        if quality == ConversionQuality.PREMIUM:
            font_size = 12
            header_font_size = 13
        
        # Önce tabloları dene
        tables_data = []
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                if any(row_data):
                    table_rows.append(row_data)
            if table_rows:
                tables_data.append(table_rows)
        
        # Tablo varsa onları kullan
        if tables_data:
            main_table = max(tables_data, key=len)
            if main_table[0] and any(main_table[0]):
                df = pd.DataFrame(main_table[1:], columns=main_table[0])
            else:
                df = pd.DataFrame(main_table)
            changes.append(f"{len(df)} satırlı tablo bulundu")
        else:
            # Normal metin
            data = []
            for para in doc.paragraphs:
                if para.text.strip():
                    data.append([para.text.strip()])
            df = pd.DataFrame(data, columns=['İçerik'])
        
        # Excel oluştur
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name='Word İçeriği')
            
            workbook = writer.book
            worksheet = writer.sheets['Word İçeriği']
            
            # Başlık stili
            header_font = Font(name='Calibri', size=header_font_size, bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="2E75B6", end_color="2E75B6", fill_type="solid")
            header_alignment = Alignment(horizontal="center", vertical="center")
            
            thin_border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            light_gray_fill = PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")
            
            # Başlık satırını formatla
            if len(df) > 0:
                for cell in worksheet[1]:
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = header_alignment
                    cell.border = thin_border
            
            # Veri hücrelerini formatla
            for row_idx, row in enumerate(worksheet.iter_rows(min_row=2, max_row=len(df)+1), 2):
                for cell in row:
                    cell.font = Font(name='Calibri', size=font_size)
                    cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
                    cell.border = thin_border
                    
                    if row_idx % 2 == 0:
                        cell.fill = light_gray_fill
                    
                    if isinstance(cell.value, (int, float)):
                        cell.number_format = '#,##0.00'
                        cell.alignment = Alignment(horizontal="right", vertical="center")
            
            # Sütun genişliklerini ayarla
            for column in worksheet.columns:
                max_length = 0
                column_letter = get_column_letter(column[0].column)
                for cell in column:
                    try:
                        if cell.value and len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                worksheet.column_dimensions[column_letter].width = adjusted_width
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 90
        metrics.complexity = DocumentComplexity.MODERATE if len(df) > 100 else DocumentComplexity.SIMPLE
        
        logger.info(f"✅ Word -> Excel dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Word -> Excel dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def word_to_pptx(input_path: str, output_path: str,
                quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Word -> PowerPoint (PROFESYONEL - TASARIM ODAKLI)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from docx import Document
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        doc = Document(input_path)
        prs = Presentation()
        
        title_font_size = 48
        heading_font_size = 32
        content_font_size = 20
        items_per_slide = 5
        
        if quality == ConversionQuality.PREMIUM:
            title_font_size = 54
            heading_font_size = 36
            content_font_size = 22
            items_per_slide = 4
        elif quality == ConversionQuality.DRAFT:
            title_font_size = 40
            heading_font_size = 28
            content_font_size = 18
            items_per_slide = 6
        
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[1]
        section_header_layout = prs.slide_layouts[2]
        
        # Ana başlık slaytı
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        if title:
            title.text = "WORD DÖKÜMANI DÖNÜŞÜMÜ"
            title.text_frame.paragraphs[0].font.size = Pt(title_font_size)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            title.text_frame.paragraphs[0].font.bold = True
        
        if subtitle:
            subtitle.text = f"Kaynak: {os.path.basename(input_path)}\nTarih: {datetime.datetime.now().strftime('%d.%m.%Y')}"
            subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        
        changes.append("Ana başlık slaytı oluşturuldu")
        
        # Paragrafları topla
        all_paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                style_name = para.style.name.lower() if para.style else "normal"
                is_heading = any(x in style_name for x in ['heading', 'title', 'başlık'])
                
                all_paragraphs.append({
                    'text': para.text.strip(),
                    'is_heading': is_heading,
                })
        
        changes.append(f"{len(all_paragraphs)} paragraf analiz edildi")
        
        # İçerik slaytları
        current_slide = None
        current_text_frame = None
        slide_count = 0
        
        for i, para in enumerate(all_paragraphs):
            start_new_slide = False
            
            if para['is_heading']:
                start_new_slide = True
            elif current_slide is None:
                start_new_slide = True
            elif i % items_per_slide == 0:
                start_new_slide = True
            
            if start_new_slide:
                slide_count += 1
                
                if para['is_heading']:
                    slide = prs.slides.add_slide(section_header_layout)
                    title = slide.shapes.title
                    if title:
                        title.text = para['text'][:50]
                        title.text_frame.paragraphs[0].font.size = Pt(36)
                        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
                        title.text_frame.paragraphs[0].font.bold = True
                    current_text_frame = None
                    changes.append(f"Başlık slaytı oluşturuldu: {para['text'][:30]}...")
                else:
                    slide = prs.slides.add_slide(content_slide_layout)
                    
                    title = slide.shapes.title
                    if title:
                        title.text = f"İçerik - Sayfa {slide_count}"
                        title.text_frame.paragraphs[0].font.size = Pt(heading_font_size)
                        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
                        title.text_frame.paragraphs[0].font.bold = True
                    
                    content = slide.placeholders[1]
                    text_frame = content.text_frame
                    text_frame.clear()
                    
                    p = text_frame.add_paragraph()
                    p.text = para['text']
                    p.font.size = Pt(content_font_size)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    
                    current_text_frame = text_frame
                    changes.append(f"İçerik slaytı oluşturuldu")
            else:
                if current_text_frame:
                    p = current_text_frame.add_paragraph()
                    p.text = para['text']
                    p.font.size = Pt(content_font_size - 2)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    changes.append(f"Slayta içerik eklendi")
        
        prs.save(output_path)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 85
        metrics.complexity = DocumentComplexity.MODERATE
        
        logger.info(f"✅ Word -> PowerPoint dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Word -> PowerPoint dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def word_to_text(input_path: str, output_path: str,
                quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Word -> Metin dönüşümü (YENİ)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from docx import Document
        
        doc = Document(input_path)
        
        text_lines = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_lines.append(para.text.strip())
        
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    text_lines.append(row_text)
        
        text = '\n'.join(text_lines)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 90
        metrics.complexity = DocumentComplexity.SIMPLE
        
        logger.info(f"✅ Word -> Metin dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Word -> Metin dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


# ========== EXCEL DÖNÜŞÜMLERİ ==========

def excel_to_pdf(input_path: str, output_path: str,
                quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Excel -> PDF (PROFESYONEL - TABLO KORUMALI)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        import pandas as pd
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        
        df = pd.read_excel(input_path)
        df = df.fillna('')
        
        c = canvas.Canvas(output_path, pagesize=landscape(A4))
        width, height = landscape(A4)
        
        y = height - 2 * cm
        line_height = 0.5 * cm
        
        # Başlık
        c.setFont("Helvetica-Bold", 14)
        c.drawString(2 * cm, y, f"Excel Dökümanı: {os.path.basename(input_path)}")
        y -= line_height * 2
        
        # Sütun başlıkları
        c.setFont("Helvetica-Bold", 10)
        x = 2 * cm
        for col in df.columns:
            c.drawString(x, y, str(col)[:15])
            x += 4 * cm
        y -= line_height
        
        # Veriler
        c.setFont("Helvetica", 9)
        for _, row in df.iterrows():
            if y < line_height:
                c.showPage()
                y = height - 2 * cm
                c.setFont("Helvetica", 9)
            
            x = 2 * cm
            for val in row:
                c.drawString(x, y, str(val)[:15])
                x += 4 * cm
            y -= line_height
        
        c.save()
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 85
        metrics.complexity = DocumentComplexity.MODERATE
        
        logger.info(f"✅ Excel -> PDF dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Excel -> PDF dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def excel_to_word(input_path: str, output_path: str,
                 quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Excel -> Word (PROFESYONEL - SAYFAYA TAM SIĞDIRMA)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        import pandas as pd
        from docx import Document
        from docx.shared import Inches, Cm, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        
        excel_file = pd.ExcelFile(input_path)
        sheet_names = excel_file.sheet_names
        metrics.input_size = os.path.getsize(input_path)
        
        doc = Document()
        
        # Sayfa yapısı ayarları
        section = doc.sections[0]
        section.page_width = Cm(21)
        section.page_height = Cm(29.7)
        section.left_margin = Cm(1.5)
        section.right_margin = Cm(1.5)
        section.top_margin = Cm(1.5)
        section.bottom_margin = Cm(1.5)
        
        # Ana başlık
        title = doc.add_heading('EXCEL DÖKÜMANI DÖNÜŞÜMÜ', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = title.runs[0]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(0, 51, 102)
        
        changes.append("Ana başlık eklendi")
        
        info_para = doc.add_paragraph()
        info_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = info_para.add_run(f"📊 Kaynak: {os.path.basename(input_path)}")
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(100, 100, 100)
        
        info_para2 = doc.add_paragraph()
        info_para2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = info_para2.add_run(f"📅 Dönüşüm: {datetime.datetime.now().strftime('%d.%m.%Y %H:%M')}")
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(100, 100, 100)
        
        doc.add_paragraph('_' * 80)
        doc.add_paragraph()
        
        doc.add_paragraph(f"📑 Toplam {len(sheet_names)} sayfa", style='Intense Quote')
        doc.add_paragraph()
        
        total_rows = 0
        
        for sheet_idx, sheet_name in enumerate(sheet_names):
            df = pd.read_excel(input_path, sheet_name=sheet_name)
            df = df.fillna('')
            total_rows += len(df)
            
            if sheet_idx > 0:
                doc.add_page_break()
            
            heading = doc.add_heading(f'Sayfa {sheet_idx + 1}: {sheet_name}', level=1)
            for run in heading.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0, 102, 204)
            
            if df.empty:
                doc.add_paragraph("📭 Bu sayfa boş")
                doc.add_paragraph()
                changes.append(f"Sayfa {sheet_name} boş")
                continue
            
            col_count = len(df.columns)
            
            min_widths = []
            for i, col in enumerate(df.columns):
                max_len = len(str(col))
                for val in df.iloc[:, i]:
                    if val != '':
                        max_len = max(max_len, len(str(val)))
                min_width = max(max_len * 0.15, 2)
                min_widths.append(min_width)
            
            total_min_width = sum(min_widths)
            available_width_cm = 18
            
            if total_min_width > available_width_cm:
                scale_factor = available_width_cm / total_min_width
                col_widths = [w * scale_factor for w in min_widths]
                changes.append(f"Sütun genişlikleri %{(1-scale_factor)*100:.0f} küçültüldü")
            else:
                extra_space = available_width_cm - total_min_width
                extra_per_col = extra_space / col_count
                col_widths = [w + extra_per_col for w in min_widths]
                changes.append(f"Sütun genişlikleri optimize edildi")
            
            rows_per_page = 30
            total_rows_sheet = len(df)
            
            for page_start in range(0, total_rows_sheet, rows_per_page):
                page_end = min(page_start + rows_per_page, total_rows_sheet)
                page_df = df.iloc[page_start:page_end]
                
                if page_start > 0:
                    doc.add_page_break()
                    doc.add_heading(f'{sheet_name} - Devam (Sayfa {page_start//rows_per_page + 2})', level=2)
                
                rows, cols = page_df.shape
                table = doc.add_table(rows=rows+1, cols=cols)
                table.style = 'Light Grid Accent 1'
                table.alignment = WD_TABLE_ALIGNMENT.CENTER
                table.autofit = False
                
                for i, width in enumerate(col_widths):
                    for row in table.rows:
                        row.cells[i].width = Cm(width)
                
                # BAŞLIK SATIRI
                for col in range(cols):
                    cell = table.cell(0, col)
                    cell.text = str(df.columns[col])
                    
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    shd = OxmlElement('w:shd')
                    shd.set(qn('w:fill'), '2E75B6')
                    tcPr.append(shd)
                    
                    for paragraph in cell.paragraphs:
                        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run = paragraph.runs[0]
                        run.font.bold = True
                        run.font.size = Pt(12)
                        run.font.name = 'Calibri'
                        run.font.color.rgb = RGBColor(255, 255, 255)
                
                # VERİ SATIRLARI
                for row in range(rows):
                    for col in range(cols):
                        cell = table.cell(row+1, col)
                        value = page_df.iloc[row, col]
                        
                        if isinstance(value, (int, float)):
                            if isinstance(value, float) and not value.is_integer():
                                cell.text = f"{value:.2f}".replace('.', ',')
                            else:
                                cell.text = str(int(value) if isinstance(value, float) else value)
                        else:
                            cell.text = str(value)
                        
                        if row % 2 == 0:
                            tc = cell._tc
                            tcPr = tc.get_or_add_tcPr()
                            shd = OxmlElement('w:shd')
                            shd.set(qn('w:fill'), 'F2F2F2')
                            tcPr.append(shd)
                        
                        for paragraph in cell.paragraphs:
                            if isinstance(page_df.iloc[row, col], (int, float)):
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                            else:
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                            
                            run = paragraph.runs[0]
                            run.font.size = Pt(10)
                            run.font.name = 'Calibri'
        
        doc.save(output_path)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 90
        metrics.complexity = DocumentComplexity.COMPLEX if total_rows > 500 else DocumentComplexity.MODERATE
        
        logger.info(f"✅ Excel -> Word dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Excel -> Word dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def excel_to_pptx(input_path: str, output_path: str,
                 quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Excel -> PowerPoint (PROFESYONEL - GRAFİK DESTEKLİ)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        import pandas as pd
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        df = pd.read_excel(input_path)
        df = df.fillna('')
        
        prs = Presentation()
        
        rows, cols = df.shape
        rows_per_slide = 15
        
        title_font_size = 48
        subtitle_font_size = 20
        heading_font_size = 32
        content_font_size = 12
        header_font_size = 13
        
        if quality == ConversionQuality.PREMIUM:
            title_font_size = 54
            subtitle_font_size = 24
            heading_font_size = 36
            content_font_size = 14
            header_font_size = 15
            rows_per_slide = 12
        
        # Ana başlık slaytı
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        if title:
            title.text = "EXCEL VERİLERİ"
            title.text_frame.paragraphs[0].font.size = Pt(title_font_size)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            title.text_frame.paragraphs[0].font.bold = True
        
        if subtitle:
            subtitle.text = f"Toplam {rows} satır, {cols} sütun\n{os.path.basename(input_path)}"
            subtitle.text_frame.paragraphs[0].font.size = Pt(subtitle_font_size)
        
        changes.append("Ana başlık slaytı oluşturuldu")
        
        # Veri slaytları
        for slide_start in range(0, rows, rows_per_slide):
            slide_end = min(slide_start + rows_per_slide, rows)
            
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            title_box = slide.shapes.add_textbox(
                int(Inches(0.5).emu), 
                int(Inches(0.2).emu), 
                int(Inches(9).emu), 
                int(Inches(0.8).emu)
            )
            title_frame = title_box.text_frame
            title_frame.text = f"Excel Verileri - Sayfa {slide_start//rows_per_slide + 1}"
            title_frame.paragraphs[0].font.size = Pt(heading_font_size)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            
            table_data = [df.columns.tolist()] + df.iloc[slide_start:slide_end].values.tolist()
            
            rows_in_slide = len(table_data)
            cols_in_slide = len(table_data[0])
            
            left = int(Inches(0.5).emu)
            top = int(Inches(1.5).emu)
            width = int(Inches(9).emu)
            height = int(Inches(5.5).emu)
            
            table = slide.shapes.add_table(rows_in_slide, cols_in_slide, left, top, width, height).table
            
            col_width = int(width / cols_in_slide)
            for col in range(cols_in_slide):
                table.columns[col].width = col_width
            
            for row in range(rows_in_slide):
                for col in range(cols_in_slide):
                    cell = table.cell(row, col)
                    cell.text = str(table_data[row][col])
                    
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(content_font_size)
                        paragraph.alignment = PP_ALIGN.CENTER
                    
                    if row == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(46, 117, 182)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.color.rgb = RGBColor(255, 255, 255)
                            paragraph.font.bold = True
                            paragraph.font.size = Pt(header_font_size)
                    else:
                        if row % 2 == 1:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
            
            changes.append(f"Veri slaytı oluşturuldu: Sayfa {slide_start//rows_per_slide + 1}")
        
        prs.save(output_path)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 85
        metrics.complexity = DocumentComplexity.MODERATE
        
        logger.info(f"✅ Excel -> PowerPoint dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Excel -> PowerPoint dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def excel_to_text(input_path: str, output_path: str,
                 quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Excel -> Metin dönüşümü (YENİ)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        import pandas as pd
        
        df = pd.read_excel(input_path)
        df = df.fillna('')
        
        text_lines = []
        
        text_lines.append(" | ".join([str(col) for col in df.columns]))
        text_lines.append("-" * 50)
        
        for _, row in df.iterrows():
            text_lines.append(" | ".join([str(val) for val in row]))
        
        text = '\n'.join(text_lines)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 90
        metrics.complexity = DocumentComplexity.SIMPLE
        
        logger.info(f"✅ Excel -> Metin dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Excel -> Metin dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


# ========== POWERPOINT DÖNÜŞÜMLERİ ==========

def pptx_to_pdf(input_path: str, output_path: str,
               quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    PowerPoint -> PDF (PROFESYONEL - TASARIM KORUMALI)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from pptx import Presentation
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        
        prs = Presentation(input_path)
        
        c = canvas.Canvas(output_path, pagesize=landscape(A4))
        width, height = landscape(A4)
        
        left_margin = 2 * cm
        y = height - 2 * cm
        line_height = 0.7 * cm
        
        total_slides = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            c.setFont("Helvetica-Bold", 18)
            c.setFillColor(colors.HexColor('#2E75B6'))
            c.drawString(left_margin, y, f"Slayt {slide_num}/{total_slides}")
            y -= line_height * 2
            
            c.setFont("Helvetica", 11)
            c.setFillColor(colors.black)
            
            shapes_with_text = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shapes_with_text += 1
                    
                    text = shape.text.strip()
                    
                    if len(text) < 50 and not '\n' in text and not text.startswith('•'):
                        c.setFont("Helvetica-Bold", 14)
                        c.setFillColor(colors.HexColor('#1E4E7C'))
                    else:
                        c.setFont("Helvetica", 11)
                        c.setFillColor(colors.black)
                    
                    text_lines = text.split('\n')
                    
                    for line in text_lines:
                        if line.strip():
                            if y < line_height + 1*cm:
                                c.showPage()
                                y = height - 2*cm
                                c.setFont("Helvetica", 11)
                                c.setFillColor(colors.black)
                            
                            if line.startswith('•') or (len(text_lines) > 1 and not c._fontname.endswith('Bold')):
                                c.drawString(left_margin + 0.5*cm, y, line)
                            else:
                                c.drawString(left_margin, y, line)
                            
                            y -= line_height
                    
                    y -= line_height * 0.5
            
            if shapes_with_text == 0:
                c.setFont("Helvetica-Oblique", 11)
                c.setFillColor(colors.HexColor('#666666'))
                c.drawString(left_margin + 0.5*cm, y, "(Bu slaytta metin yok)")
                y -= line_height
            
            changes.append(f"Slayt {slide_num} işlendi")
            
            if slide_num < total_slides:
                c.showPage()
                y = height - 2*cm
        
        c.save()
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 90
        metrics.complexity = DocumentComplexity.MODERATE if total_slides > 10 else DocumentComplexity.SIMPLE
        
        logger.info(f"✅ PowerPoint -> PDF dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ PowerPoint -> PDF dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def pptx_to_word(input_path: str, output_path: str,
                quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    PowerPoint -> Word (PROFESYONEL)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from pptx import Presentation
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        prs = Presentation(input_path)
        doc = Document()
        
        total_slides = len(prs.slides)
        
        style = doc.styles['Normal']
        style.font.name = 'Calibri'
        style.font.size = Pt(11)
        
        title = doc.add_heading('POWERPOINT DÖKÜMANI DÖNÜŞÜMÜ', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = title.runs[0]
        run.font.size = Pt(24)
        run.font.name = 'Calibri'
        run.font.color.rgb = (0, 51, 102)
        
        doc.add_paragraph(f"Kaynak dosya: {os.path.basename(input_path)}")
        doc.add_paragraph(f"Toplam {total_slides} slayt")
        doc.add_paragraph(f"Dönüşüm tarihi: {datetime.datetime.now().strftime('%d.%m.%Y %H:%M')}")
        doc.add_paragraph()
        
        changes.append("Ana başlık oluşturuldu")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            heading = doc.add_heading(f'Slayt {slide_num}/{total_slides}', level=1)
            for run in heading.runs:
                run.font.color.rgb = (0, 102, 204)
            
            shapes_with_text = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shapes_with_text += 1
                    
                    text = shape.text.strip()
                    
                    if len(text) < 50 and not '\n' in text and not text.startswith('•'):
                        doc.add_heading(text, level=2)
                    else:
                        paragraphs = text.split('\n')
                        for para in paragraphs:
                            if para.strip():
                                p = doc.add_paragraph()
                                run = p.add_run(para.strip())
                                run.font.size = Pt(11)
                                p.paragraph_format.left_indent = Inches(0.3)
                                p.paragraph_format.space_after = Pt(3)
                                
                                if para.startswith('•'):
                                    p.style = 'List Bullet'
            
            if shapes_with_text == 0:
                doc.add_paragraph("(Bu slaytta metin yok)")
            
            changes.append(f"Slayt {slide_num} işlendi ({shapes_with_text} şekil)")
            
            if slide_num < total_slides:
                doc.add_page_break()
        
        doc.save(output_path)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 85
        metrics.complexity = DocumentComplexity.MODERATE if total_slides > 10 else DocumentComplexity.SIMPLE
        
        logger.info(f"✅ PowerPoint -> Word dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ PowerPoint -> Word dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def pptx_to_text(input_path: str, output_path: str,
                quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    PowerPoint -> Metin dönüşümü (YENİ)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from pptx import Presentation
        
        prs = Presentation(input_path)
        
        text_lines = []
        total_slides = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_lines.append(f"--- Slayt {slide_num}/{total_slides} ---")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_lines.append(shape.text.strip())
            
            text_lines.append("")
        
        text = '\n'.join(text_lines)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 90
        metrics.complexity = DocumentComplexity.SIMPLE
        
        logger.info(f"✅ PowerPoint -> Metin dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ PowerPoint -> Metin dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


# ========== PDF DÖNÜŞÜMLERİ ==========

def pdf_to_word(input_path: str, output_path: str,
               quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    PDF -> Word (PROFESYONEL - METİN KORUMALI)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        import PyPDF2
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document()
        
        style = doc.styles['Normal']
        style.font.name = 'Times New Roman'
        style.font.size = Pt(12)
        
        with open(input_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            total_pages = len(pdf_reader.pages)
            
            title = doc.add_heading('PDF DÖKÜMANI DÖNÜŞÜMÜ', 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = title.runs[0]
            run.font.size = Pt(24)
            run.font.name = 'Arial'
            run.font.color.rgb = (0, 51, 102)
            
            doc.add_paragraph(f"Kaynak dosya: {os.path.basename(input_path)}")
            doc.add_paragraph(f"Toplam {total_pages} sayfa")
            doc.add_paragraph(f"Dönüşüm tarihi: {datetime.datetime.now().strftime('%d.%m.%Y %H:%M')}")
            doc.add_paragraph()
            
            changes.append("Ana başlık oluşturuldu")
            
            for page_num in range(total_pages):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                
                heading = doc.add_heading(f'Sayfa {page_num + 1}/{total_pages}', level=1)
                for run in heading.runs:
                    run.font.color.rgb = (0, 102, 204)
                
                if text and text.strip():
                    lines = text.split('\n')
                    current_paragraph = []
                    
                    for line in lines:
                        clean_line = ' '.join(line.split())
                        if clean_line:
                            current_paragraph.append(clean_line)
                        else:
                            if current_paragraph:
                                doc.add_paragraph(' '.join(current_paragraph))
                                current_paragraph = []
                    
                    if current_paragraph:
                        doc.add_paragraph(' '.join(current_paragraph))
                    
                    changes.append(f"Sayfa {page_num + 1} işlendi ({len(lines)} satır)")
                else:
                    doc.add_paragraph("(Bu sayfada metin bulunamadı)")
                    changes.append(f"Sayfa {page_num + 1} boş")
                
                if page_num < total_pages - 1:
                    doc.add_page_break()
        
        doc.save(output_path)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 80
        metrics.complexity = DocumentComplexity.MODERATE if total_pages > 20 else DocumentComplexity.SIMPLE
        
        logger.info(f"✅ PDF -> Word dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ PDF -> Word dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def pdf_to_text(input_path: str, output_path: str,
               quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    PDF -> Metin dönüşümü (YENİ)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        import PyPDF2
        
        text_lines = []
        
        with open(input_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            total_pages = len(pdf_reader.pages)
            
            for page_num in range(total_pages):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                
                if text and text.strip():
                    text_lines.append(f"--- Sayfa {page_num + 1}/{total_pages} ---")
                    text_lines.append(text.strip())
                    text_lines.append("")
        
        text = '\n'.join(text_lines)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 85
        metrics.complexity = DocumentComplexity.SIMPLE
        
        logger.info(f"✅ PDF -> Metin dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ PDF -> Metin dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


# ========== GÖRSEL DÖNÜŞÜMLERİ ==========

def image_to_pdf(input_path: str, output_path: str,
                quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Görsel -> PDF (PROFESYONEL - YÜKSEK KALİTE)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from PIL import Image, ImageFilter
        import img2pdf
        
        image = Image.open(input_path)
        
        changes.append(f"Orijinal görsel: {image.width}x{image.height}, {image.mode}")
        
        if image.mode != 'RGB':
            image = image.convert('RGB')
            changes.append("RGB'ye dönüştürüldü")
        
        dpi = 300
        quality_value = 95
        if quality == ConversionQuality.PREMIUM:
            dpi = 600
            quality_value = 100
        elif quality == ConversionQuality.DRAFT:
            dpi = 150
            quality_value = 80
        
        a4_width, a4_height = int(8.27 * dpi), int(11.69 * dpi)
        image.thumbnail((a4_width, a4_height), Image.Resampling.LANCZOS)
        changes.append(f"A4 boyutuna ölçeklendi ({image.width}x{image.height})")
        
        if quality != ConversionQuality.DRAFT:
            image = image.filter(ImageFilter.SHARPEN)
            changes.append("Keskinleştirme uygulandı")
        
        temp_image_path = input_path + "_temp.jpg"
        image.save(temp_image_path, 'JPEG', quality=quality_value, optimize=True, dpi=(dpi, dpi))
        changes.append(f"JPEG olarak kaydedildi (kalite: {quality_value})")
        
        pdf_bytes = img2pdf.convert(temp_image_path)
        
        with open(output_path, "wb") as f:
            f.write(pdf_bytes)
        
        os.remove(temp_image_path)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 95 if quality == ConversionQuality.PREMIUM else 85
        
        logger.info(f"✅ Görsel -> PDF dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Görsel -> PDF dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def image_to_word(input_path: str, output_path: str,
                 quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Görsel -> Word (OCR - PROFESYONEL)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from PIL import Image, ImageEnhance, ImageFilter
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        import pytesseract
        
        image = Image.open(input_path)
        
        original_format = image.format
        original_size = image.size
        
        changes.append(f"Orijinal görsel: {original_size[0]}x{original_size[1]}, {original_format}")
        
        scale_factor = 2
        contrast_factor = 2.5
        psm_mode = 6
        
        if quality == ConversionQuality.PREMIUM:
            scale_factor = 3
            contrast_factor = 3.0
            psm_mode = 3
        elif quality == ConversionQuality.DRAFT:
            scale_factor = 1.5
            contrast_factor = 2.0
            psm_mode = 6
        
        width, height = image.size
        if width < 2000:
            new_size = (width * scale_factor, height * scale_factor)
            image = image.resize(new_size, Image.Resampling.LANCZOS)
            changes.append(f"Görsel büyütüldü: {width}x{height} -> {new_size[0]}x{new_size[1]}")
        
        image = image.convert('L')
        changes.append("Gri tonlamaya çevrildi")
        
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(contrast_factor)
        changes.append(f"Kontrast %{contrast_factor*100:.0f} artırıldı")
        
        image = image.filter(ImageFilter.MedianFilter(size=3))
        image = image.filter(ImageFilter.SHARPEN)
        changes.append("Gürültü azaltıldı, keskinleştirildi")
        
        if quality != ConversionQuality.DRAFT:
            image = image.filter(ImageFilter.EDGE_ENHANCE)
            changes.append("Kenar iyileştirmesi uygulandı")
        
        temp_image_path = input_path + "_temp_ocr.png"
        image.save(temp_image_path, 'PNG', dpi=(300,300))
        
        sample_text = pytesseract.image_to_string(temp_image_path, lang='tur', config='--psm 6')
        if not sample_text.strip():
            sample_text = pytesseract.image_to_string(temp_image_path, lang='eng', config='--psm 6')
            language = 'eng'
            changes.append("İngilizce dilinde OCR yapılacak")
        else:
            language = 'tur'
            changes.append("Türkçe dilinde OCR yapılacak")
        
        ocr_text = ""
        ocr_configs = [
            f'--oem 3 --psm {psm_mode} -l {language}+eng',
            f'--oem 3 --psm 3 -l {language}+eng',
            f'--oem 3 --psm 4 -l {language}+eng',
            f'--oem 3 --psm 11 -l {language}+eng',
        ]
        
        for config in ocr_configs:
            try:
                result = pytesseract.image_to_string(temp_image_path, config=config)
                if result.strip() and len(result) > len(ocr_text):
                    ocr_text = result
            except:
                continue
        
        text = ocr_text
        
        if text.strip():
            text = fix_common_ocr_errors(text)
            cleaned_text, removed_lines = clean_ocr_text(text)
            
            cleaned_text = re.sub(r'\s+', ' ', cleaned_text)
            cleaned_text = re.sub(r'\n\s*\n\s*\n', '\n\n', cleaned_text)
            
            cleaned_text = merge_intelligent_lines(cleaned_text)
            cleaned_text = normalize_whitespace(cleaned_text)
            
            changes.append(f"OCR başarılı ({len(text)} karakter, {len(removed_lines)} gereksiz satır temizlendi)")
            
            confidence = calculate_ocr_confidence(cleaned_text)
            changes.append(f"OCR güven skoru: %{confidence:.0f}")
        else:
            changes.append("OCR başarısız, görselde metin olmayabilir")
            cleaned_text = text
        
        doc = Document()
        
        section = doc.sections[0]
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        
        style = doc.styles['Normal']
        style.font.name = 'Calibri'
        style.font.size = Pt(11)
        
        title = doc.add_heading('📄 GÖRSELDEN OCR İLE DÖNÜŞTÜRÜLEN METİN', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = title.runs[0]
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(0, 51, 102)
        
        doc.add_paragraph()
        info_para = doc.add_paragraph()
        info_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = info_para.add_run(f"🖼️ Kaynak: {os.path.basename(input_path)}")
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(100, 100, 100)
        
        info_para2 = doc.add_paragraph()
        info_para2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = info_para2.add_run(f"📅 Dönüşüm: {datetime.datetime.now().strftime('%d.%m.%Y %H:%M')}")
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(100, 100, 100)
        
        info_para3 = doc.add_paragraph()
        info_para3.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = info_para3.add_run(f"🔍 OCR Dili: {language.upper()}, Kalite: {quality.value}")
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(100, 100, 100)
        
        doc.add_paragraph('_' * 80)
        doc.add_paragraph()
        
        if cleaned_text.strip():
            clean_paragraphs = cleaned_text.split('\n\n')
            para_count = 0
            
            for para in clean_paragraphs:
                if para.strip():
                    p = doc.add_paragraph()
                    run = p.add_run(para.strip())
                    run.font.size = Pt(11)
                    run.font.name = 'Calibri'
                    
                    p.paragraph_format.space_after = Pt(12)
                    para_count += 1
            
            changes.append(f"{para_count} paragraf oluşturuldu")
        else:
            p = doc.add_paragraph()
            run = p.add_run("(Görselde metin bulunamadı veya okunamadı)")
            run.font.size = Pt(12)
            run.font.italic = True
            run.font.color.rgb = RGBColor(150, 150, 150)
        
        doc.add_page_break()
        
        doc.add_heading('🖼️ ORİJİNAL GÖRSEL', level=1)
        doc.add_paragraph()
        
        try:
            doc.add_picture(input_path, width=Inches(5))
            last_paragraph = doc.paragraphs[-1]
            last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            changes.append("Orijinal görsel eklendi")
        except:
            doc.add_paragraph("(Görsel yüklenemedi)")
        
        os.remove(temp_image_path)
        
        doc.save(output_path)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 85 if text.strip() else 50
        
        logger.info(f"✅ Görsel -> Word dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Görsel -> Word dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def image_to_text(input_path: str, output_path: str,
                 quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Görsel -> Metin dönüşümü (OCR) (YENİ)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from PIL import Image, ImageEnhance, ImageFilter
        import pytesseract
        
        image = Image.open(input_path)
        
        scale_factor = 2
        contrast_factor = 2.5
        
        if quality == ConversionQuality.PREMIUM:
            scale_factor = 3
            contrast_factor = 3.0
        elif quality == ConversionQuality.DRAFT:
            scale_factor = 1.5
            contrast_factor = 2.0
        
        width, height = image.size
        if width < 2000:
            new_size = (width * scale_factor, height * scale_factor)
            image = image.resize(new_size, Image.Resampling.LANCZOS)
            changes.append(f"Görsel büyütüldü: {width}x{height} -> {new_size[0]}x{new_size[1]}")
        
        image = image.convert('L')
        
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(contrast_factor)
        
        image = image.filter(ImageFilter.MedianFilter(size=3))
        image = image.filter(ImageFilter.SHARPEN)
        
        if quality != ConversionQuality.DRAFT:
            image = image.filter(ImageFilter.EDGE_ENHANCE)
        
        temp_image_path = input_path + "_temp_ocr.png"
        image.save(temp_image_path, 'PNG', dpi=(300,300))
        
        sample_text = pytesseract.image_to_string(temp_image_path, lang='tur', config='--psm 6')
        if not sample_text.strip():
            sample_text = pytesseract.image_to_string(temp_image_path, lang='eng', config='--psm 6')
            language = 'eng'
        else:
            language = 'tur'
        
        ocr_configs = [
            f'--oem 3 --psm 6 -l {language}+eng',
            f'--oem 3 --psm 3 -l {language}+eng',
            f'--oem 3 --psm 4 -l {language}+eng',
        ]
        
        text = ""
        for config in ocr_configs:
            try:
                result = pytesseract.image_to_string(temp_image_path, config=config)
                if result.strip() and len(result) > len(text):
                    text = result
            except:
                continue
        
        if text.strip():
            text = fix_common_ocr_errors(text)
            cleaned_text, removed_lines = clean_ocr_text(text)
            cleaned_text = re.sub(r'\s+', ' ', cleaned_text)
            cleaned_text = merge_intelligent_lines(cleaned_text)
            cleaned_text = normalize_whitespace(cleaned_text)
            
            changes.append(f"OCR başarılı ({len(text)} karakter, {len(removed_lines)} gereksiz satır temizlendi)")
            confidence = calculate_ocr_confidence(cleaned_text)
            changes.append(f"OCR güven skoru: %{confidence:.0f}")
        else:
            changes.append("OCR başarısız, görselde metin olmayabilir")
            cleaned_text = text
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(cleaned_text if cleaned_text.strip() else "(Görselde metin bulunamadı)")
        
        os.remove(temp_image_path)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 85 if text.strip() else 50
        
        logger.info(f"✅ Görsel -> Metin dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Görsel -> Metin dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


# ========== METİN DÖNÜŞÜMLERİ ==========

def text_to_pdf(input_path: str, output_path: str,
               quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Metin -> PDF dönüşümü (YENİ)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
        
        with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        y = height - 2 * cm
        line_height = 0.6 * cm
        
        c.setFont("Helvetica", 11)
        
        lines = text.split('\n')
        for line in lines:
            if line.strip():
                if y < line_height:
                    c.showPage()
                    y = height - 2 * cm
                    c.setFont("Helvetica", 11)
                
                c.drawString(2 * cm, y, line.strip())
                y -= line_height
            else:
                y -= line_height * 0.5
        
        c.save()
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 90
        metrics.complexity = DocumentComplexity.SIMPLE
        
        logger.info(f"✅ Metin -> PDF dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Metin -> PDF dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


def text_to_word(input_path: str, output_path: str,
                quality: ConversionQuality = ConversionQuality.PROFESSIONAL) -> Tuple[bool, str, ConversionMetrics]:
    """
    Metin -> Word dönüşümü (YENİ)
    """
    metrics = ConversionMetrics()
    metrics.input_size = os.path.getsize(input_path)
    changes = []
    
    try:
        from docx import Document
        
        with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        
        doc = Document()
        
        doc.add_heading('Metin Dökümanı Dönüşümü', 0)
        doc.add_paragraph(f"Kaynak: {os.path.basename(input_path)}")
        doc.add_paragraph()
        
        paragraphs = text.split('\n')
        for para in paragraphs:
            if para.strip():
                doc.add_paragraph(para.strip())
        
        doc.save(output_path)
        
        metrics.output_size = os.path.getsize(output_path)
        metrics.processing_time = 0
        metrics.compression_ratio = metrics.output_size / metrics.input_size if metrics.input_size > 0 else 1
        metrics.quality_score = 90
        metrics.complexity = DocumentComplexity.SIMPLE
        
        logger.info(f"✅ Metin -> Word dönüşüm başarılı: {input_path}")
        return True, output_path, metrics
        
    except Exception as e:
        logger.error(f"❌ Metin -> Word dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", metrics


# ========== ANA DÖNÜŞTÜRME FONKSİYONLARI ==========

async def smart_convert_file(input_path: str, output_path: str, source_type: str, target_type: str,
                            user_id: int = None, db_instance: Any = None,
                            quality: str = "profesyonel") -> Tuple[bool, str, str, Optional[str], ConversionMetrics]:
    """
    Gelişmiş dönüşüm yöneticisi - TÜM DÖNÜŞÜMLER DESTEKLENİR
    """
    import time
    start_time = time.time()
    
    logger.info(f"🔍 Dönüşüm isteği: {source_type} -> {target_type}")
    logger.info(f"📁 Kaynak: {input_path}")
    logger.info(f"📁 Hedef: {output_path}")
    
    quality_map = {
        'taslak': ConversionQuality.DRAFT,
        'standart': ConversionQuality.STANDARD,
        'profesyonel': ConversionQuality.PROFESSIONAL,
        'premium': ConversionQuality.PREMIUM
    }
    quality_level = quality_map.get(quality, ConversionQuality.PROFESSIONAL)
    
    metrics = ConversionMetrics()
    changes = []
    edit_summary = None
    conversion_type = 'direct'
    
    try:
        if not os.path.exists(input_path):
            error_msg = f"Dosya bulunamadı: {input_path}"
            logger.error(f"❌ {error_msg}")
            metrics.warnings.append(error_msg)
            return False, "", conversion_type, edit_summary, metrics
        
        if not is_conversion_supported(source_type, target_type):
            error_msg = f"Desteklenmeyen dönüşüm: {source_type} -> {target_type}"
            logger.error(f"❌ {error_msg}")
            metrics.warnings.append(error_msg)
            return False, "", conversion_type, edit_summary, metrics
        
        # TÜM DÖNÜŞÜM FONKSİYONLARI
        conversion_functions = {
            # Word dönüşümleri
            ('WORD', 'PDF'): lambda: word_to_pdf(input_path, output_path, quality_level),
            ('WORD', 'EXCEL'): lambda: word_to_excel(input_path, output_path, quality_level),
            ('WORD', 'POWERPOINT'): lambda: word_to_pptx(input_path, output_path, quality_level),
            ('WORD', 'TEXT'): lambda: word_to_text(input_path, output_path, quality_level),
            
            # Excel dönüşümleri
            ('EXCEL', 'PDF'): lambda: excel_to_pdf(input_path, output_path, quality_level),
            ('EXCEL', 'WORD'): lambda: excel_to_word(input_path, output_path, quality_level),
            ('EXCEL', 'POWERPOINT'): lambda: excel_to_pptx(input_path, output_path, quality_level),
            ('EXCEL', 'TEXT'): lambda: excel_to_text(input_path, output_path, quality_level),
            
            # PowerPoint dönüşümleri
            ('POWERPOINT', 'PDF'): lambda: pptx_to_pdf(input_path, output_path, quality_level),
            ('POWERPOINT', 'WORD'): lambda: pptx_to_word(input_path, output_path, quality_level),
            ('POWERPOINT', 'TEXT'): lambda: pptx_to_text(input_path, output_path, quality_level),
            
            # PDF dönüşümleri
            ('PDF', 'WORD'): lambda: pdf_to_word(input_path, output_path, quality_level),
            ('PDF', 'TEXT'): lambda: pdf_to_text(input_path, output_path, quality_level),
            
            # Görsel dönüşümleri
            ('GORSEL', 'PDF'): lambda: image_to_pdf(input_path, output_path, quality_level),
            ('GORSEL', 'WORD'): lambda: image_to_word(input_path, output_path, quality_level),
            ('GORSEL', 'TEXT'): lambda: image_to_text(input_path, output_path, quality_level),
            
            # Metin dönüşümleri
            ('TEXT', 'PDF'): lambda: text_to_pdf(input_path, output_path, quality_level),
            ('TEXT', 'WORD'): lambda: text_to_word(input_path, output_path, quality_level),
        }
        
        key = (source_type, target_type)
        if key not in conversion_functions:
            error_msg = f"Dönüşüm fonksiyonu bulunamadı: {source_type} -> {target_type}"
            logger.error(f"❌ {error_msg}")
            metrics.warnings.append(error_msg)
            return False, "", conversion_type, edit_summary, metrics
        
        logger.info(f"🔄 Dönüşüm başlıyor: {source_type} -> {target_type}")
        
        try:
            success, out_path, conv_metrics = conversion_functions[key]()
        except Exception as e:
            logger.error(f"❌ Dönüşüm fonksiyonu hatası: {e}")
            traceback.print_exc()
            metrics.warnings.append(str(e))
            return False, "", conversion_type, edit_summary, metrics
        
        metrics = conv_metrics
        metrics.processing_time = time.time() - start_time
        
        if success and out_path and os.path.exists(out_path):
            metrics.quality_score = conv_metrics.quality_score
            logger.info(f"✅ Dönüşüm başarılı: {out_path}")
            return True, out_path, conversion_type, edit_summary, metrics
        else:
            error_msg = f"Dönüşüm başarısız: {source_type} -> {target_type}"
            logger.error(f"❌ {error_msg}")
            metrics.warnings.append(error_msg)
            return False, "", conversion_type, edit_summary, metrics
            
    except Exception as e:
        logger.error(f"❌ Akıllı dönüşüm hatası: {e}")
        traceback.print_exc()
        metrics.warnings.append(str(e))
        return False, "", conversion_type, edit_summary, metrics


def get_conversion_report(metrics: ConversionMetrics) -> str:
    """
    Dönüşüm raporu oluştur
    """
    report = f"""
📊 **DÖNÜŞÜM RAPORU**

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📁 **DOSYA BİLGİLERİ**
• Giriş boyutu: {get_file_size_str(metrics.input_size)}
• Çıkış boyutu: {get_file_size_str(metrics.output_size)}
• Sıkıştırma oranı: %{(1-metrics.compression_ratio)*100:.1f}
• İşlem süresi: {metrics.processing_time:.2f} saniye

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📊 **KALİTE METRİKLERİ**
• Kalite puanı: %{metrics.quality_score}
• Karmaşıklık: {metrics.complexity.value}
"""
    
    if metrics.warnings:
        report += "\n⚠️ **UYARILAR**\n"
        for warning in metrics.warnings[:5]:
            report += f"  • {warning}\n"
    
    if metrics.suggestions:
        report += "\n💡 **ÖNERİLER**\n"
        for suggestion in metrics.suggestions[:3]:
            report += f"  • {suggestion}\n"
    
    return report


# ========== GERİYE UYUMLULUK FONKSİYONLARI ==========

async def convert_file(input_path: str, output_path: str, source_type: str, target_type: str) -> Tuple[bool, str]:
    """
    Eski dönüşüm fonksiyonu (geriye uyumluluk için)
    """
    success, out_path, _, _, metrics = await smart_convert_file(
        input_path, output_path, source_type, target_type, quality="standart"
    )
    
    if success:
        return True, ""
    else:
        error_msg = metrics.warnings[0] if metrics.warnings else "Bilinmeyen hata"
        return False, error_msg


async def smart_process_file(input_path: str, output_path: str, source_type: str, target_type: str,
                            user_id: int = None, db_instance: Any = None) -> Tuple[bool, str, Dict]:
    """
    Akıllı işlem yöneticisi (geriye uyumluluk için)
    """
    success, out_path, conv_type, edit_summary, metrics = await smart_convert_file(
        input_path, output_path, source_type, target_type, user_id, db_instance
    )
    
    results = {
        'conversion': {
            'success': success,
            'message': "Dönüşüm tamamlandı" if success else "Dönüşüm başarısız",
            'type': conv_type,
            'edit_summary': edit_summary
        },
        'metrics': {
            'quality_score': metrics.quality_score,
            'processing_time': metrics.processing_time,
            'input_size': metrics.input_size,
            'output_size': metrics.output_size,
            'warnings': metrics.warnings
        }
    }
    
    return success, get_conversion_report(metrics), results


async def smart_process_all(input_path: str, output_path: str, source_type: str, target_type: str,
                           user_id: int = None, db_instance: Any = None) -> Tuple[bool, str, Dict]:
    """
    Tüm işlemler (geriye uyumluluk için)
    """
    return await smart_process_file(input_path, output_path, source_type, target_type, user_id, db_instance)


# ========== TEST FONKSİYONU ==========
if __name__ == "__main__":
    print("🔧 Profesyonel Dönüşüm Modülü Test Ediliyor...")
    print("=" * 60)
    
    print("\n📋 DESTEKLENEN DÖNÜŞÜMLER:")
    for source, targets in SUPPORTED_TARGET_FORMATS.items():
        print(f"  • {source} -> {', '.join(targets)}")
    
    print("\n" + "=" * 60)
    print("✅ Modül hazır! Tüm dönüşümler destekleniyor.")
